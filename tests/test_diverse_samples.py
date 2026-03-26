"""GroundedDeck 多样化样本数据集回归测试。

覆盖：
- saas-launch 英文 source pack 的完整管线
- 所有 fixture 的渲染兼容性
- 布局多样性验证
- 英文关键词规则引擎增强
"""

from __future__ import annotations

import json
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from src.renderer.pptx_renderer import render_slide_spec_to_pptx

ROOT = Path(__file__).resolve().parent.parent
SAAS_SPEC = ROOT / "fixtures" / "slide-spec" / "saas-launch-slide-spec.json"
SAAS_PACK = ROOT / "fixtures" / "source-packs" / "saas-launch-source-pack.json"
TECH_SPEC = ROOT / "fixtures" / "slide-spec" / "tech-review-slide-spec.json"
TECH_PACK = ROOT / "fixtures" / "source-packs" / "tech-review-source-pack.json"
TECH_SPEC = ROOT / "fixtures" / "slide-spec" / "tech-review-slide-spec.json"
TECH_PACK = ROOT / "fixtures" / "source-packs" / "tech-review-source-pack.json"


def load_json(p: Path) -> dict:
    return json.loads(p.read_text(encoding="utf-8"))


# ── SaaS Launch 渲染测试 ──────────────────────────────────────────────


class SaasLaunchRenderingTests(unittest.TestCase):
    """saas-launch 英文样本渲染测试。"""

    def test_renders_successfully(self) -> None:
        spec = load_json(SAAS_SPEC)
        with tempfile.TemporaryDirectory() as tmpdir:
            out = render_slide_spec_to_pptx(spec, Path(tmpdir) / "t.pptx")
            self.assertTrue(out.exists())
            self.assertGreater(out.stat().st_size, 10_000)

    def test_correct_slide_count(self) -> None:
        spec = load_json(SAAS_SPEC)
        with tempfile.TemporaryDirectory() as tmpdir:
            out = render_slide_spec_to_pptx(spec, Path(tmpdir) / "t.pptx")
            prs = Presentation(str(out))
            self.assertEqual(len(prs.slides), 7)

    def test_diverse_layouts(self) -> None:
        spec = load_json(SAAS_SPEC)
        layouts = {s["layout_type"] for s in spec["slides"]}
        self.assertGreaterEqual(len(layouts), 3, f"Expected diverse layouts, got: {layouts}")

    def test_has_timeline(self) -> None:
        spec = load_json(SAAS_SPEC)
        self.assertTrue(any(s["layout_type"] == "timeline" for s in spec["slides"]))

    def test_has_comparison(self) -> None:
        spec = load_json(SAAS_SPEC)
        self.assertTrue(any(s["layout_type"] == "comparison" for s in spec["slides"]))

    def test_has_process(self) -> None:
        spec = load_json(SAAS_SPEC)
        self.assertTrue(any(s["layout_type"] == "process" for s in spec["slides"]))

    def test_has_chart(self) -> None:
        spec = load_json(SAAS_SPEC)
        self.assertTrue(any(s["layout_type"] == "chart" for s in spec["slides"]))

    def test_english_content_in_pptx(self) -> None:
        spec = load_json(SAAS_SPEC)
        with tempfile.TemporaryDirectory() as tmpdir:
            out = render_slide_spec_to_pptx(spec, Path(tmpdir) / "t.pptx")
            prs = Presentation(str(out))
            all_text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        all_text += shape.text_frame.text
            self.assertIn("AnalyticsPro", all_text)


# ── SaaS Launch 管线测试 ──────────────────────────────────────────────


class SaasLaunchPipelineTests(unittest.TestCase):
    """saas-launch 完整管线测试。"""

    def test_quality_passes(self) -> None:
        from src.ingest.normalize import normalize_source_pack
        from src.llm.provider import DeterministicProvider
        from src.quality.checks import grade_slide_spec

        raw = load_json(SAAS_PACK)
        norm = normalize_source_pack(raw)
        provider = DeterministicProvider()
        spec = provider.draft_slide_spec(norm)
        qr = grade_slide_spec(norm, spec, provider=provider)
        self.assertEqual(qr["status"], "pass")

    def test_pipeline_end_to_end(self) -> None:
        from src.runtime.pipeline import run_pipeline

        raw = load_json(SAAS_PACK)
        with tempfile.TemporaryDirectory() as tmpdir:
            result = run_pipeline(raw, render_pptx=Path(tmpdir) / "t.pptx")
            self.assertIn("slide_spec", result)
            self.assertIn("quality_report", result)
            self.assertIn("pptx_path", result)
            self.assertEqual(result["quality_report"]["status"], "pass")
            self.assertIn("artifact_grade", result)

    def test_narrative_grade(self) -> None:
        from src.runtime.pipeline import run_pipeline

        raw = load_json(SAAS_PACK)
        with tempfile.TemporaryDirectory() as tmpdir:
            result = run_pipeline(raw, render_pptx=Path(tmpdir) / "t.pptx")
            self.assertIn("narrative_grade", result)
            self.assertIn("status", result["narrative_grade"])


# ── 所有 Fixture 兼容性测试 ───────────────────────────────────────────


class AllFixturesCompatibilityTests(unittest.TestCase):
    """所有 fixture 的渲染和质量兼容性测试。"""

    def test_all_fixtures_render(self) -> None:
        fixture_dir = ROOT / "fixtures" / "slide-spec"
        for spec_file in sorted(fixture_dir.glob("*.json")):
            with self.subTest(fixture=spec_file.name):
                spec = load_json(spec_file)
                with tempfile.TemporaryDirectory() as tmpdir:
                    out = render_slide_spec_to_pptx(spec, Path(tmpdir) / "t.pptx")
                    self.assertTrue(out.exists())
                    self.assertGreater(out.stat().st_size, 5_000)

    def test_all_source_packs_quality(self) -> None:
        from src.ingest.normalize import normalize_source_pack
        from src.llm.provider import DeterministicProvider
        from src.quality.checks import grade_slide_spec

        fixture_dir = ROOT / "fixtures" / "source-packs"
        provider = DeterministicProvider()
        for pack_file in sorted(fixture_dir.glob("*.json")):
            with self.subTest(fixture=pack_file.name):
                raw = load_json(pack_file)
                norm = normalize_source_pack(raw)
                spec = provider.draft_slide_spec(norm)
                qr = grade_slide_spec(norm, spec, provider=provider)
                self.assertEqual(
                    qr["status"], "pass",
                    f"{pack_file.name} quality failed: {qr.get('failures')}",
                )


# ── 英文关键词规则引擎测试 ────────────────────────────────────────────


class EnglishKeywordInferenceTests(unittest.TestCase):
    """英文关键词规则引擎增强测试。"""

    def test_landscape_infers_comparison(self) -> None:
        from src.visual.selector import infer_layout_type

        unit = {
            "section_heading": "Competitive Landscape",
            "text": "Competitive Landscape Incumbent tools dominate enterprise but leave mid-market underserved.",
        }
        self.assertEqual(infer_layout_type(unit).layout_type, "comparison")

    def test_phase_pattern_infers_process(self) -> None:
        from src.visual.selector import infer_layout_type

        unit = {
            "section_heading": "Launch Sequence",
            "text": "Launch Sequence Phase 1: beta launch. Phase 2: public launch. Phase 3: channel partners.",
        }
        self.assertEqual(infer_layout_type(unit).layout_type, "process")

    def test_sequence_keyword_infers_process(self) -> None:
        from src.visual.selector import infer_layout_type

        unit = {
            "section_heading": "Go-to-Market Sequence",
            "text": "Go-to-Market Sequence First build partnerships then expand to new regions.",
        }
        self.assertEqual(infer_layout_type(unit).layout_type, "process")

    def test_grew_from_with_years_infers_timeline(self) -> None:
        from src.visual.selector import infer_layout_type

        unit = {
            "section_heading": "Market Size Timeline",
            "text": "Market Size Timeline The segment grew from $2.1B in 2022 to $3.4B in 2024.",
        }
        self.assertEqual(infer_layout_type(unit).layout_type, "timeline")

    def test_english_phase_step_extraction(self) -> None:
        from src.visual.selector import _extract_process_steps

        text = "Phase 1: invite-only beta with 50 partners. Phase 2: public launch. Phase 3: channel enablement."
        steps = _extract_process_steps(text)
        self.assertGreaterEqual(len(steps), 2)
        self.assertIn("invite-only beta with 50 partners", steps[0])

    def test_chinese_comparison_regression(self) -> None:
        """中文 comparison 关键词应仍然正常工作。"""
        from src.visual.selector import infer_layout_type

        unit = {
            "section_heading": "欧洲 vs 东南亚对比",
            "text": "欧洲 vs 东南亚对比 欧洲单车收入更高，但认证压力更重。",
        }
        self.assertEqual(infer_layout_type(unit).layout_type, "comparison")

    def test_chinese_process_regression(self) -> None:
        """中文 process 关键词应仍然正常工作。"""
        from src.visual.selector import infer_layout_type

        unit = {
            "section_heading": "进入路径",
            "text": "进入路径 建议先以东南亚市场建立经销，再把验证过的运营模型迁移到欧洲。",
        }
        self.assertEqual(infer_layout_type(unit).layout_type, "process")


# ── Tech Review 渲染测试 ──────────────────────────────────────────────


class TechReviewRenderingTests(unittest.TestCase):
    """tech-review 英文样本渲染测试。"""

    def test_renders_successfully(self) -> None:
        spec = load_json(TECH_SPEC)
        with tempfile.TemporaryDirectory() as tmpdir:
            out = render_slide_spec_to_pptx(spec, Path(tmpdir) / "t.pptx")
            self.assertTrue(out.exists())
            self.assertGreater(out.stat().st_size, 10_000)

    def test_correct_slide_count(self) -> None:
        spec = load_json(TECH_SPEC)
        with tempfile.TemporaryDirectory() as tmpdir:
            out = render_slide_spec_to_pptx(spec, Path(tmpdir) / "t.pptx")
            prs = Presentation(str(out))
            self.assertEqual(len(prs.slides), 11)

    def test_seven_distinct_layouts(self) -> None:
        spec = load_json(TECH_SPEC)
        layouts = {s["layout_type"] for s in spec["slides"]}
        self.assertEqual(len(layouts), 7, f"Expected 7 distinct layouts, got: {layouts}")

    def test_has_section_dividers(self) -> None:
        spec = load_json(TECH_SPEC)
        section_slides = [s for s in spec["slides"] if s["layout_type"] == "section"]
        self.assertEqual(len(section_slides), 3)

    def test_has_all_content_layouts(self) -> None:
        spec = load_json(TECH_SPEC)
        layouts = {s["layout_type"] for s in spec["slides"]}
        for expected in ("timeline", "comparison", "process", "chart", "section"):
            self.assertIn(expected, layouts, f"Missing layout: {expected}")

    def test_heading_boost_timeline(self) -> None:
        """标题含 'Timeline' + 年份应优先匹配 timeline 而非 comparison。"""
        from src.visual.selector import infer_layout_type

        unit = {
            "section_heading": "System Evolution Timeline",
            "text": "System Evolution Timeline CloudStack was originally deployed in 2019. In 2021 the team extracted auth. By 2023 the API gateway was added but the core remains monolithic.",
        }
        self.assertEqual(infer_layout_type(unit).layout_type, "timeline")

    def test_checkpoint_heading_infers_timeline(self) -> None:
        """标题含 'Checkpoint' + 年份应匹配 timeline。"""
        from src.visual.selector import infer_layout_type

        unit = {
            "section_heading": "Risk Timeline and Checkpoints",
            "text": "Risk Timeline and Checkpoints Risk exposure peaks during Q1 2026 when the order pipeline runs in dual-write mode.",
        }
        self.assertEqual(infer_layout_type(unit).layout_type, "timeline")

    def test_monolith_vs_microservices_infers_comparison(self) -> None:
        """'Monolith vs Microservices' 应匹配 comparison。"""
        from src.visual.selector import infer_layout_type

        unit = {
            "section_heading": "Monolith vs Microservices Trade-offs",
            "text": "Monolith vs Microservices Trade-offs The monolithic order engine offers simpler debugging but limits independent scaling.",
        }
        self.assertEqual(infer_layout_type(unit).layout_type, "comparison")

    def test_latency_throughput_infers_chart(self) -> None:
        """含 latency/throughput 关键词应匹配 chart。"""
        from src.visual.selector import infer_layout_type

        unit = {
            "section_heading": "Latency and Throughput Metrics",
            "text": "Latency and Throughput Metrics Current P99 latency is 420ms with peak throughput of 2800 TPS.",
        }
        self.assertEqual(infer_layout_type(unit).layout_type, "chart")


# ── Section 分隔页测试 ────────────────────────────────────────────────


class SectionDividerTests(unittest.TestCase):
    """section 分隔页自动插入测试。"""

    def test_no_section_dividers_for_small_packs(self) -> None:
        """4 个 content slides 不应插入 section 分隔页。"""
        from src.runtime.pipeline import run_pipeline

        raw = load_json(ROOT / "fixtures" / "source-packs" / "strongest-demo-source-pack.json")
        result = run_pipeline(raw)
        layouts = [s["layout_type"] for s in result["slide_spec"]["slides"]]
        self.assertNotIn("section", layouts)

    def test_no_section_dividers_for_5_units(self) -> None:
        """5 个 content slides 不应插入 section 分隔页。"""
        from src.runtime.pipeline import run_pipeline

        raw = load_json(SAAS_PACK)
        result = run_pipeline(raw)
        layouts = [s["layout_type"] for s in result["slide_spec"]["slides"]]
        self.assertNotIn("section", layouts)

    def test_section_dividers_for_6_plus_units(self) -> None:
        """6+ 个 content slides 应自动插入 section 分隔页。"""
        from src.runtime.pipeline import run_pipeline

        raw = load_json(TECH_PACK)
        result = run_pipeline(raw)
        layouts = [s["layout_type"] for s in result["slide_spec"]["slides"]]
        self.assertIn("section", layouts)

    def test_section_dividers_between_different_sources(self) -> None:
        """section 分隔页应出现在不同 source 之间。"""
        from src.runtime.pipeline import run_pipeline

        raw = load_json(TECH_PACK)
        result = run_pipeline(raw)
        slides = result["slide_spec"]["slides"]
        section_titles = [s["title"] for s in slides if s["layout_type"] == "section"]
        self.assertEqual(len(section_titles), 3)

    def test_section_quality_passes(self) -> None:
        """含 section 分隔页的 deck 应通过质量检查。"""
        from src.runtime.pipeline import run_pipeline

        raw = load_json(TECH_PACK)
        result = run_pipeline(raw)
        self.assertEqual(result["quality_report"]["status"], "pass")

    def test_tech_review_pipeline_end_to_end(self) -> None:
        """tech-review 端到端管线测试。"""
        from src.runtime.pipeline import run_pipeline

        raw = load_json(TECH_PACK)
        with tempfile.TemporaryDirectory() as tmpdir:
            result = run_pipeline(raw, render_pptx=Path(tmpdir) / "t.pptx")
            self.assertIn("pptx_path", result)
            self.assertEqual(result["quality_report"]["status"], "pass")
            self.assertIn("artifact_grade", result)
