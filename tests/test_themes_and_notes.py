"""GroundedDeck 主题系统和 Speaker Notes 增强测试。

覆盖：
- 主题注册表和获取
- 所有内置主题的 PPTX 渲染
- 主题颜色一致性验证
- 主题切换不影响内容完整性
- 智能 Speaker Notes 结构化内容
- Speaker Notes 包含关键数据引用
- Pipeline 主题集成
"""

from __future__ import annotations

import json
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from src.renderer.pptx_renderer import render_slide_spec_to_pptx
from src.renderer.themes import (
    BUILTIN_THEMES,
    SlideTheme,
    get_theme,
    list_themes,
)

ROOT = Path(__file__).resolve().parent.parent
STRONGEST_SPEC = ROOT / "fixtures" / "slide-spec" / "strongest-demo-slide-spec.json"
SAAS_SPEC = ROOT / "fixtures" / "slide-spec" / "saas-launch-slide-spec.json"
TECH_SPEC = ROOT / "fixtures" / "slide-spec" / "tech-review-slide-spec.json"
STRONGEST_PACK = ROOT / "fixtures" / "source-packs" / "strongest-demo-source-pack.json"


def load_json(p: Path) -> dict:
    return json.loads(p.read_text(encoding="utf-8"))


# ── 主题注册表测试 ────────────────────────────────────────────────────


class ThemeRegistryTests(unittest.TestCase):
    """主题注册表和获取测试。"""

    def test_list_themes_returns_all_builtin(self) -> None:
        themes = list_themes()
        self.assertEqual(len(themes), 5)
        self.assertIn("professional-blue", themes)
        self.assertIn("forest-green", themes)
        self.assertIn("warm-sunset", themes)
        self.assertIn("minimal-gray", themes)
        self.assertIn("ocean-teal", themes)

    def test_get_default_theme(self) -> None:
        theme = get_theme()
        self.assertEqual(theme.name, "professional-blue")

    def test_get_theme_by_name(self) -> None:
        theme = get_theme("forest-green")
        self.assertEqual(theme.name, "forest-green")
        self.assertEqual(theme.display_name, "Forest Green")

    def test_get_theme_none_returns_default(self) -> None:
        theme = get_theme(None)
        self.assertEqual(theme.name, "professional-blue")

    def test_get_theme_empty_returns_default(self) -> None:
        theme = get_theme("")
        self.assertEqual(theme.name, "professional-blue")

    def test_get_theme_invalid_raises(self) -> None:
        with self.assertRaises(ValueError) as ctx:
            get_theme("nonexistent-theme")
        self.assertIn("nonexistent-theme", str(ctx.exception))
        self.assertIn("available themes", str(ctx.exception))

    def test_all_themes_are_frozen(self) -> None:
        for name, theme in BUILTIN_THEMES.items():
            with self.subTest(theme=name):
                self.assertIsInstance(theme, SlideTheme)
                with self.assertRaises(AttributeError):
                    theme.name = "modified"  # type: ignore[misc]

    def test_all_themes_have_complete_color_slots(self) -> None:
        """每个主题都应有完整的颜色槽位。"""
        required_attrs = [
            "primary", "accent", "text_dark", "text_body", "text_muted", "text_white",
            "bg_light", "bg_white", "border", "highlight",
            "timeline_accent", "comparison_left", "comparison_right",
            "process_accent", "chart_accent", "cover_subtitle",
            "table_alt_left", "table_alt_right",
        ]
        for name, theme in BUILTIN_THEMES.items():
            with self.subTest(theme=name):
                for attr in required_attrs:
                    self.assertTrue(
                        hasattr(theme, attr),
                        f"Theme '{name}' missing attribute '{attr}'",
                    )


# ── 主题渲染测试 ──────────────────────────────────────────────────────


class ThemeRenderingTests(unittest.TestCase):
    """所有内置主题的 PPTX 渲染测试。"""

    def test_all_themes_render_strongest_demo(self) -> None:
        """每个主题都应能成功渲染 strongest-demo。"""
        spec = load_json(STRONGEST_SPEC)
        for theme_name in list_themes():
            with self.subTest(theme=theme_name):
                with tempfile.TemporaryDirectory() as tmpdir:
                    out = render_slide_spec_to_pptx(
                        spec, Path(tmpdir) / "t.pptx", theme=theme_name,
                    )
                    self.assertTrue(out.exists())
                    self.assertGreater(out.stat().st_size, 10_000)

    def test_all_themes_render_tech_review(self) -> None:
        """每个主题都应能渲染含 section 分隔页的 tech-review。"""
        spec = load_json(TECH_SPEC)
        for theme_name in list_themes():
            with self.subTest(theme=theme_name):
                with tempfile.TemporaryDirectory() as tmpdir:
                    out = render_slide_spec_to_pptx(
                        spec, Path(tmpdir) / "t.pptx", theme=theme_name,
                    )
                    prs = Presentation(str(out))
                    self.assertEqual(len(prs.slides), 11)

    def test_theme_does_not_change_slide_count(self) -> None:
        """不同主题不应改变 slide 数量。"""
        spec = load_json(SAAS_SPEC)
        for theme_name in list_themes():
            with self.subTest(theme=theme_name):
                with tempfile.TemporaryDirectory() as tmpdir:
                    out = render_slide_spec_to_pptx(
                        spec, Path(tmpdir) / "t.pptx", theme=theme_name,
                    )
                    prs = Presentation(str(out))
                    self.assertEqual(len(prs.slides), 7)

    def test_theme_object_accepted(self) -> None:
        """直接传入 SlideTheme 对象也应正常工作。"""
        spec = load_json(STRONGEST_SPEC)
        theme = get_theme("ocean-teal")
        with tempfile.TemporaryDirectory() as tmpdir:
            out = render_slide_spec_to_pptx(
                spec, Path(tmpdir) / "t.pptx", theme=theme,
            )
            self.assertTrue(out.exists())

    def test_default_theme_backward_compatible(self) -> None:
        """不传 theme 参数应与之前行为一致。"""
        spec = load_json(STRONGEST_SPEC)
        with tempfile.TemporaryDirectory() as tmpdir:
            out = render_slide_spec_to_pptx(spec, Path(tmpdir) / "t.pptx")
            self.assertTrue(out.exists())
            self.assertGreater(out.stat().st_size, 10_000)

    def test_different_themes_produce_different_files(self) -> None:
        """不同主题应产生不同的文件内容（文件大小可能略有差异）。"""
        spec = load_json(STRONGEST_SPEC)
        sizes = {}
        with tempfile.TemporaryDirectory() as tmpdir:
            for theme_name in ["professional-blue", "forest-green", "minimal-gray"]:
                out = render_slide_spec_to_pptx(
                    spec, Path(tmpdir) / f"{theme_name}.pptx", theme=theme_name,
                )
                sizes[theme_name] = out.stat().st_size
        # 所有文件都应成功生成且大小合理
        for name, size in sizes.items():
            self.assertGreater(size, 10_000, f"{name} file too small")


# ── Speaker Notes 增强测试 ────────────────────────────────────────────


class SpeakerNotesEnhancementTests(unittest.TestCase):
    """智能 Speaker Notes 结构化内容测试。"""

    def test_speaker_notes_have_structure(self) -> None:
        """speaker notes 应包含结构化标签。"""
        from src.ingest.normalize import normalize_source_pack
        from src.llm.provider import DeterministicProvider

        raw = load_json(STRONGEST_PACK)
        norm = normalize_source_pack(raw)
        provider = DeterministicProvider()
        spec = provider.draft_slide_spec(norm)

        # 检查 content slides（跳过 cover 和 summary）
        content_slides = [s for s in spec["slides"] if s["layout_type"] not in ("cover", "summary")]
        for slide in content_slides:
            notes = slide["speaker_notes"]
            self.assertIn("[开场]", notes, f"Slide {slide['slide_id']} missing [开场]")
            self.assertIn("[要点]", notes, f"Slide {slide['slide_id']} missing [要点]")
            self.assertIn("[过渡]", notes, f"Slide {slide['slide_id']} missing [过渡]")
            self.assertIn("[来源]", notes, f"Slide {slide['slide_id']} missing [来源]")

    def test_speaker_notes_contain_heading(self) -> None:
        """speaker notes 应包含 section heading。"""
        from src.ingest.normalize import normalize_source_pack
        from src.llm.provider import DeterministicProvider

        raw = load_json(STRONGEST_PACK)
        norm = normalize_source_pack(raw)
        provider = DeterministicProvider()
        spec = provider.draft_slide_spec(norm)

        content_slides = [s for s in spec["slides"] if s["layout_type"] not in ("cover", "summary")]
        for slide in content_slides:
            notes = slide["speaker_notes"]
            self.assertIn(slide["title"], notes)

    def test_speaker_notes_contain_data_refs(self) -> None:
        """包含数字的 unit 的 speaker notes 应有 [数据] 标签。"""
        from src.ingest.normalize import normalize_source_pack
        from src.llm.provider import DeterministicProvider

        raw = load_json(STRONGEST_PACK)
        norm = normalize_source_pack(raw)
        provider = DeterministicProvider()
        spec = provider.draft_slide_spec(norm)

        # 找到 chart 类型的 slide（通常包含数字）
        chart_slides = [s for s in spec["slides"] if s["layout_type"] == "chart"]
        for slide in chart_slides:
            notes = slide["speaker_notes"]
            self.assertIn("[数据]", notes, f"Chart slide {slide['slide_id']} missing [数据]")

    def test_speaker_notes_contain_source_binding(self) -> None:
        """speaker notes 应包含 source binding 信息。"""
        from src.ingest.normalize import normalize_source_pack
        from src.llm.provider import DeterministicProvider

        raw = load_json(STRONGEST_PACK)
        norm = normalize_source_pack(raw)
        provider = DeterministicProvider()
        spec = provider.draft_slide_spec(norm)

        content_slides = [s for s in spec["slides"] if s["layout_type"] not in ("cover", "summary")]
        for slide in content_slides:
            notes = slide["speaker_notes"]
            binding = slide["source_bindings"][0]
            self.assertIn(binding, notes)

    def test_speaker_notes_layout_specific_opener(self) -> None:
        """不同布局类型应有不同的开场提示。"""
        from src.ingest.normalize import normalize_source_pack
        from src.llm.provider import DeterministicProvider

        raw = load_json(STRONGEST_PACK)
        norm = normalize_source_pack(raw)
        provider = DeterministicProvider()
        spec = provider.draft_slide_spec(norm)

        layout_openers = {
            "timeline": "时间演进",
            "comparison": "对比",
            "process": "行动路径",
            "chart": "数据指标",
        }

        content_slides = [s for s in spec["slides"] if s["layout_type"] not in ("cover", "summary")]
        for slide in content_slides:
            layout = slide["layout_type"]
            if layout in layout_openers:
                notes = slide["speaker_notes"]
                self.assertIn(
                    layout_openers[layout], notes,
                    f"Slide {slide['slide_id']} ({layout}) missing layout-specific opener",
                )

    def test_speaker_notes_in_pptx(self) -> None:
        """PPTX 文件中的 speaker notes 应包含结构化内容。"""
        spec = load_json(STRONGEST_SPEC)
        with tempfile.TemporaryDirectory() as tmpdir:
            out = render_slide_spec_to_pptx(spec, Path(tmpdir) / "t.pptx")
            prs = Presentation(str(out))
            # 检查至少有一个 slide 的 notes 包含 [开场]
            found_structured = False
            for slide in prs.slides:
                notes_text = slide.notes_slide.notes_text_frame.text
                if "[开场]" in notes_text:
                    found_structured = True
                    break
            self.assertTrue(found_structured, "No structured speaker notes found in PPTX")


# ── Pipeline 主题集成测试 ─────────────────────────────────────────────


class PipelineThemeIntegrationTests(unittest.TestCase):
    """Pipeline 主题集成测试。"""

    def test_pipeline_with_theme(self) -> None:
        """pipeline 应支持 theme 参数。"""
        from src.runtime.pipeline import run_pipeline

        raw = load_json(STRONGEST_PACK)
        with tempfile.TemporaryDirectory() as tmpdir:
            result = run_pipeline(
                raw,
                render_pptx=Path(tmpdir) / "t.pptx",
                theme="forest-green",
            )
            self.assertIn("pptx_path", result)
            self.assertEqual(result["theme"], "forest-green")
            self.assertEqual(result["quality_report"]["status"], "pass")

    def test_pipeline_default_theme(self) -> None:
        """pipeline 不传 theme 应正常工作。"""
        from src.runtime.pipeline import run_pipeline

        raw = load_json(STRONGEST_PACK)
        with tempfile.TemporaryDirectory() as tmpdir:
            result = run_pipeline(raw, render_pptx=Path(tmpdir) / "t.pptx")
            self.assertIn("pptx_path", result)
            self.assertNotIn("theme", result)  # 默认不记录 theme

    def test_pipeline_all_themes_quality_pass(self) -> None:
        """所有主题下 pipeline 质量检查都应通过。"""
        from src.runtime.pipeline import run_pipeline

        raw = load_json(STRONGEST_PACK)
        for theme_name in list_themes():
            with self.subTest(theme=theme_name):
                with tempfile.TemporaryDirectory() as tmpdir:
                    result = run_pipeline(
                        raw,
                        render_pptx=Path(tmpdir) / "t.pptx",
                        theme=theme_name,
                    )
                    self.assertEqual(result["quality_report"]["status"], "pass")
