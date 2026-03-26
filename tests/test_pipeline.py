from __future__ import annotations

import json
import tempfile
import unittest
from pathlib import Path

from src.ingest.normalize import normalize_source_pack
from src.llm.provider import (
    DeterministicProvider,
    OpenAICompatibleProvider,
    ProviderConfig,
    build_provider_from_env,
)
from src.planner.draft import draft_slide_spec
from src.quality.checks import grade_slide_spec
from src.runtime.cli import write_pipeline_outputs, write_verification_summary
from src.runtime.demo import write_strongest_demo_bundle
from src.runtime.pipeline import run_pipeline


ROOT = Path(__file__).resolve().parent.parent
RAW_FIXTURE = ROOT / "fixtures" / "source-packs" / "example-source-pack.json"
NORMALIZED_FIXTURE = ROOT / "fixtures" / "normalized-source-units" / "example-normalized-source-units.json"
SLIDE_SPEC_FIXTURE = ROOT / "fixtures" / "slide-spec" / "example-slide-spec.json"
STRONGEST_DEMO_RAW_FIXTURE = ROOT / "fixtures" / "source-packs" / "strongest-demo-source-pack.json"
STRONGEST_DEMO_NORMALIZED_FIXTURE = (
    ROOT / "fixtures" / "normalized-source-units" / "strongest-demo-normalized-source-units.json"
)
STRONGEST_DEMO_SLIDE_SPEC_FIXTURE = ROOT / "fixtures" / "slide-spec" / "strongest-demo-slide-spec.json"
STRONGEST_DEMO_QUALITY_FIXTURE = ROOT / "fixtures" / "quality-reports" / "strongest-demo-quality-report.json"
SLIDE_SCHEMA = ROOT / "schemas" / "slide-spec.schema.json"


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


class PipelineFixtureTests(unittest.TestCase):
    def test_fixtures_exist(self) -> None:
        self.assertTrue(RAW_FIXTURE.exists(), RAW_FIXTURE)
        self.assertTrue(NORMALIZED_FIXTURE.exists(), NORMALIZED_FIXTURE)
        self.assertTrue(SLIDE_SPEC_FIXTURE.exists(), SLIDE_SPEC_FIXTURE)
        self.assertTrue(STRONGEST_DEMO_RAW_FIXTURE.exists(), STRONGEST_DEMO_RAW_FIXTURE)
        self.assertTrue(STRONGEST_DEMO_NORMALIZED_FIXTURE.exists(), STRONGEST_DEMO_NORMALIZED_FIXTURE)
        self.assertTrue(STRONGEST_DEMO_SLIDE_SPEC_FIXTURE.exists(), STRONGEST_DEMO_SLIDE_SPEC_FIXTURE)
        self.assertTrue(STRONGEST_DEMO_QUALITY_FIXTURE.exists(), STRONGEST_DEMO_QUALITY_FIXTURE)

    def test_normalize_source_pack_matches_fixture(self) -> None:
        raw_pack = load_json(RAW_FIXTURE)
        normalized = normalize_source_pack(raw_pack)
        expected = load_json(NORMALIZED_FIXTURE)
        self.assertEqual(normalized, expected)

    def test_draft_slide_spec_matches_fixture(self) -> None:
        normalized = load_json(NORMALIZED_FIXTURE)
        slide_spec = draft_slide_spec(normalized)
        expected = load_json(SLIDE_SPEC_FIXTURE)
        self.assertEqual(slide_spec, expected)

    def test_drafted_slide_spec_passes_quality_checks(self) -> None:
        normalized = load_json(NORMALIZED_FIXTURE)
        slide_spec = draft_slide_spec(normalized)
        report = grade_slide_spec(normalized, slide_spec)
        self.assertEqual(report["status"], "pass")
        self.assertFalse(report["failures"])
        self.assertEqual(report["grounding"]["grounding_ratio"], 1.0)
        self.assertEqual(report["visual_form"]["match_ratio"], 1.0)

    def test_explicit_deterministic_provider_keeps_fixture_output_stable(self) -> None:
        normalized = load_json(NORMALIZED_FIXTURE)
        provider = DeterministicProvider()

        slide_spec = draft_slide_spec(normalized, provider=provider)
        report = grade_slide_spec(normalized, slide_spec, provider=provider)

        self.assertEqual(slide_spec, load_json(SLIDE_SPEC_FIXTURE))
        self.assertEqual(report["status"], "pass")
        self.assertEqual(report["provider"], "deterministic")

    def test_strongest_demo_normalize_source_pack_matches_fixture(self) -> None:
        raw_pack = load_json(STRONGEST_DEMO_RAW_FIXTURE)
        normalized = normalize_source_pack(raw_pack)
        expected = load_json(STRONGEST_DEMO_NORMALIZED_FIXTURE)
        self.assertEqual(normalized, expected)

    def test_strongest_demo_draft_slide_spec_matches_fixture(self) -> None:
        normalized = load_json(STRONGEST_DEMO_NORMALIZED_FIXTURE)
        slide_spec = draft_slide_spec(normalized)
        expected = load_json(STRONGEST_DEMO_SLIDE_SPEC_FIXTURE)
        self.assertEqual(slide_spec, expected)

    def test_strongest_demo_quality_report_matches_fixture(self) -> None:
        normalized = load_json(STRONGEST_DEMO_NORMALIZED_FIXTURE)
        slide_spec = load_json(STRONGEST_DEMO_SLIDE_SPEC_FIXTURE)
        report = grade_slide_spec(normalized, slide_spec)
        expected = load_json(STRONGEST_DEMO_QUALITY_FIXTURE)
        self.assertEqual(report, expected)

    def test_strongest_demo_uses_diverse_visual_forms(self) -> None:
        slide_spec = load_json(STRONGEST_DEMO_SLIDE_SPEC_FIXTURE)
        layout_types = {slide["layout_type"] for slide in slide_spec["slides"]}
        self.assertTrue({"timeline", "comparison", "process", "chart"}.issubset(layout_types))

    def test_build_provider_from_env_defaults_to_deterministic(self) -> None:
        provider = build_provider_from_env({})
        self.assertIsInstance(provider, DeterministicProvider)

    def test_provider_config_round_trip(self) -> None:
        provider = DeterministicProvider(
            ProviderConfig(
                provider="deterministic",
                model="baseline-fixture",
                api_key_env="IGNORED",
                base_url="http://localhost:11434",
            )
        )

        self.assertEqual(provider.config.provider, "deterministic")
        self.assertEqual(provider.config.model, "baseline-fixture")
        self.assertEqual(provider.config.base_url, "http://localhost:11434")

    def test_build_provider_from_env_supports_openai_compatible(self) -> None:
        provider = build_provider_from_env(
            {
                "GROUNDED_DECK_LLM_PROVIDER": "openai-compatible",
                "GROUNDED_DECK_LLM_MODEL": "gpt-4.1-mini",
                "GROUNDED_DECK_BASE_URL": "https://api.example.com/v1",
                "GROUNDED_DECK_API_KEY": "secret",
            }
        )

        self.assertIsInstance(provider, OpenAICompatibleProvider)
        self.assertEqual(provider.config.model, "gpt-4.1-mini")
        self.assertEqual(provider.config.base_url, "https://api.example.com/v1")

    def test_openai_compatible_provider_requires_base_url(self) -> None:
        with self.assertRaises(ValueError):
            build_provider_from_env(
                {
                    "GROUNDED_DECK_LLM_PROVIDER": "openai-compatible",
                    "GROUNDED_DECK_LLM_MODEL": "gpt-4.1-mini",
                    "GROUNDED_DECK_API_KEY": "secret",
                }
            )

    def test_openai_compatible_provider_builds_chat_request(self) -> None:
        provider = OpenAICompatibleProvider(
            ProviderConfig(
                provider="openai-compatible",
                model="gpt-4.1-mini",
                api_key_env="GROUNDED_DECK_API_KEY",
                base_url="https://api.example.com/v1",
            ),
            api_key="secret",
        )

        request = provider.build_chat_request(
            system_prompt="You are a planner.",
            user_prompt="Summarize the source pack.",
            response_format={"type": "json_object"},
        )

        self.assertEqual(request["url"], "https://api.example.com/v1/chat/completions")
        self.assertEqual(request["headers"]["Authorization"], "Bearer secret")
        self.assertEqual(request["json"]["model"], "gpt-4.1-mini")
        self.assertEqual(request["json"]["response_format"], {"type": "json_object"})
        self.assertEqual(len(request["json"]["messages"]), 2)

    def test_openai_compatible_provider_adds_reasoning_split_for_minimax(self) -> None:
        provider = OpenAICompatibleProvider(
            ProviderConfig(
                provider="openai-compatible",
                model="MiniMax-M2.7",
                api_key_env="GROUNDED_DECK_API_KEY",
                base_url="https://api.minimaxi.com/v1",
            ),
            api_key="secret",
        )

        request = provider.build_chat_request(
            system_prompt="You are a planner.",
            user_prompt="Summarize the source pack.",
            response_format={"type": "json_object"},
        )

        self.assertTrue(request["json"]["reasoning_split"])
        self.assertEqual(request["json"]["response_format"], {"type": "json_object"})

    def test_openai_compatible_provider_parses_json_content(self) -> None:
        provider = OpenAICompatibleProvider(
            ProviderConfig(
                provider="openai-compatible",
                model="gpt-4.1-mini",
                api_key_env="GROUNDED_DECK_API_KEY",
                base_url="https://api.example.com/v1",
            ),
            api_key="secret",
        )

        parsed = provider.parse_json_response(
            {
                "choices": [
                    {
                        "message": {
                            "content": '{"status":"pass","failures":[],"coverage":{"required_units":3,"covered_units":3},"grounding":{"total_content_slides":2,"grounded_slides":2},"visual_form":{"expected_units":3,"matched_units":3}}'
                        }
                    }
                ]
            }
        )

        self.assertEqual(parsed["status"], "pass")
        self.assertEqual(parsed["coverage"]["covered_units"], 3)

    def test_openai_compatible_provider_parses_json_content_wrapped_in_code_fence(self) -> None:
        provider = OpenAICompatibleProvider(
            ProviderConfig(
                provider="openai-compatible",
                model="gpt-4.1-mini",
                api_key_env="GROUNDED_DECK_API_KEY",
                base_url="https://api.example.com/v1",
            ),
            api_key="secret",
        )

        parsed = provider.parse_json_response(
            {
                "choices": [
                    {
                        "message": {
                            "content": "```json\n{\"status\":\"pass\",\"failures\":[],\"coverage\":{\"required_units\":3,\"covered_units\":3},\"grounding\":{\"total_content_slides\":2,\"grounded_slides\":2},\"visual_form\":{\"expected_units\":3,\"matched_units\":3}}\n```"
                        }
                    }
                ]
            }
        )

        self.assertEqual(parsed["status"], "pass")
        self.assertEqual(parsed["visual_form"]["matched_units"], 3)

    def test_strongest_demo_planner_prompt_includes_archived_acceptance_guardrails(self) -> None:
        normalized = load_json(STRONGEST_DEMO_NORMALIZED_FIXTURE)

        prompt = OpenAICompatibleProvider.build_planner_user_prompt(normalized)

        self.assertIn("Strongest-demo accepted live baseline:", prompt)
        self.assertIn("Produce exactly 6 slides in this order", prompt)
        self.assertIn("Use the exact slide_id sequence ['slide-01', 'slide-02', 'slide-03', 'slide-04', 'slide-05', 'slide-06'].", prompt)
        self.assertIn("China EV Market Entry: Europe & Southeast Asia Strategy", prompt)
        self.assertIn("src-01:sec-01->timeline", prompt)
        self.assertIn("must use source_bindings ['src-01:sec-01'] and must_include_checks ['src-01:sec-01']", prompt)
        self.assertIn("source_bindings set to [] and must_include_checks set to []", prompt)
        self.assertIn("Decision Backbone', use layout_type 'summary', set source_bindings to all source units, and set must_include_checks to exactly []", prompt)
        self.assertIn("Slide 6 source_bindings must be exactly ['src-01:sec-01', 'src-01:sec-02', 'src-02:sec-01', 'src-03:sec-01'] in that order.", prompt)
        self.assertIn("generated_at_unix as archival metadata only", prompt)
        self.assertIn("grounded_content_slides must be 5 of 5", prompt)
        self.assertIn("covered_unit_ids must be exactly ['src-01:sec-01', 'src-01:sec-02', 'src-02:sec-01', 'src-03:sec-01']", prompt)
        self.assertIn("visual_matched_unit_ids must be exactly ['src-01:sec-01', 'src-01:sec-02', 'src-02:sec-01', 'src-03:sec-01']", prompt)

    def test_strongest_demo_grader_prompt_includes_archived_acceptance_checks(self) -> None:
        normalized = load_json(STRONGEST_DEMO_NORMALIZED_FIXTURE)
        slide_spec = load_json(STRONGEST_DEMO_SLIDE_SPEC_FIXTURE)

        prompt = OpenAICompatibleProvider.build_grader_user_prompt(normalized, slide_spec)

        self.assertIn("Strongest-demo accepted live baseline checks:", prompt)
        self.assertIn("Compare the output structurally against the archived acceptance snapshot", prompt)
        self.assertIn("strongest-demo-1774370225/acceptance-summary.json", prompt)
        self.assertIn("generated_at_unix as the only tolerated archival delta", prompt)
        self.assertIn("Fail if slide_count is not exactly 6.", prompt)
        self.assertIn("Fail if slide_id sequence is not exactly ['slide-01', 'slide-02', 'slide-03', 'slide-04', 'slide-05', 'slide-06'].", prompt)
        self.assertIn("Fail if layout_sequence is not exactly", prompt)
        self.assertIn("China EV Market Entry: Europe & Southeast Asia Strategy", prompt)
        self.assertIn("source_bindings == [] and must_include_checks == []", prompt)
        self.assertIn("Fail if any strongest-demo unit slide does not keep source_bindings == [unit_id] and must_include_checks == [unit_id].", prompt)
        self.assertIn("Fail if the final slide is not 'Decision Backbone' with layout_type 'summary', source_bindings exactly ['src-01:sec-01', 'src-01:sec-02', 'src-02:sec-01', 'src-03:sec-01'], and must_include_checks == [].", prompt)
        self.assertIn("Fail if strongest-demo bilingual unit slide titles drift", prompt)
        self.assertIn("Fail if grounded_content_slides is not 5 or total_content_slides is not 5.", prompt)
        self.assertIn("Fail if covered_unit_ids is not exactly ['src-01:sec-01', 'src-01:sec-02', 'src-02:sec-01', 'src-03:sec-01'].", prompt)
        self.assertIn("Fail if visual_matched_unit_ids is not exactly ['src-01:sec-01', 'src-01:sec-02', 'src-02:sec-01', 'src-03:sec-01'].", prompt)

    def test_openai_compatible_provider_parses_text_content_list(self) -> None:
        provider = OpenAICompatibleProvider(
            ProviderConfig(
                provider="openai-compatible",
                model="gpt-4.1-mini",
                api_key_env="GROUNDED_DECK_API_KEY",
                base_url="https://api.example.com/v1",
            ),
            api_key="secret",
        )

        parsed = provider.parse_json_response(
            {
                "choices": [
                    {
                        "message": {
                            "content": [
                                {
                                    "type": "text",
                                    "text": '{"status":"pass","failures":[],"coverage":{"required_units":3,"covered_units":3},"grounding":{"total_content_slides":2,"grounded_slides":2},"visual_form":{"expected_units":3,"matched_units":3}}',
                                }
                            ]
                        }
                    }
                ]
            }
        )

        self.assertEqual(parsed["status"], "pass")
        self.assertEqual(parsed["grounding"]["grounded_slides"], 2)

    def test_openai_compatible_provider_parses_minimax_think_wrapped_json(self) -> None:
        provider = OpenAICompatibleProvider(
            ProviderConfig(
                provider="openai-compatible",
                model="MiniMax-M2.7",
                api_key_env="GROUNDED_DECK_API_KEY",
                base_url="https://api.minimaxi.com/v1",
            ),
            api_key="secret",
        )

        parsed = provider.parse_json_response(
            {
                "choices": [
                    {
                        "message": {
                            "content": "<think>internal reasoning with example {\\\"ok\\\": false}</think>\n\n{\"status\":\"pass\",\"failures\":[],\"coverage\":{\"required_units\":3,\"covered_units\":3},\"grounding\":{\"total_content_slides\":2,\"grounded_slides\":2},\"visual_form\":{\"expected_units\":3,\"matched_units\":3}}"
                        }
                    }
                ]
            }
        )

        self.assertEqual(parsed["status"], "pass")
        self.assertEqual(parsed["coverage"]["required_units"], 3)

    def test_openai_compatible_provider_reports_content_snippet_on_invalid_json(self) -> None:
        provider = OpenAICompatibleProvider(
            ProviderConfig(
                provider="openai-compatible",
                model="gpt-4.1-mini",
                api_key_env="GROUNDED_DECK_API_KEY",
                base_url="https://api.example.com/v1",
            ),
            api_key="secret",
        )

        with self.assertRaisesRegex(ValueError, "content_snippet=.*not json"):
            provider.parse_json_response(
                {
                    "choices": [
                        {
                            "message": {
                                "content": "not json at all"
                            }
                        }
                    ]
                }
            )

    def test_openai_compatible_provider_reports_transport_body_snippet_on_invalid_json(self) -> None:
        with self.assertRaisesRegex(ValueError, "body_snippet=.*not json"):
            OpenAICompatibleProvider._parse_transport_body("not json from upstream")

    def test_openai_compatible_provider_can_draft_with_mock_transport(self) -> None:
        expected = load_json(SLIDE_SPEC_FIXTURE)
        normalized = load_json(STRONGEST_DEMO_NORMALIZED_FIXTURE)

        def fake_transport(request: dict) -> dict:
            self.assertEqual(request["json"]["response_format"], {"type": "json_object"})
            self.assertIn("Return only valid JSON matching the slide-spec schema", request["json"]["messages"][0]["content"])
            self.assertIn("Then add exactly one content slide per source unit.", request["json"]["messages"][1]["content"])
            self.assertIn("must_include_checks to exactly [unit_id]", request["json"]["messages"][1]["content"])
            self.assertIn("timeline for chronology/year-based evidence", request["json"]["messages"][1]["content"])
            self.assertIn("china-ev-market-entry", request["json"]["messages"][1]["content"])
            return {
                "choices": [
                    {
                        "message": {
                            "content": json.dumps(expected, ensure_ascii=False)
                        }
                    }
                ]
            }

        provider = OpenAICompatibleProvider(
            ProviderConfig(
                provider="openai-compatible",
                model="gpt-4.1-mini",
                api_key_env="GROUNDED_DECK_API_KEY",
                base_url="https://api.example.com/v1",
            ),
            api_key="secret",
            transport=fake_transport,
        )

        drafted = provider.draft_slide_spec(normalized)

        # 验证 _layout_validation 附加字段存在且结构正确
        self.assertIn("_layout_validation", drafted)
        validation = drafted["_layout_validation"]
        self.assertIn("matched", validation)
        self.assertIn("mismatched", validation)
        self.assertIn("total", validation)
        self.assertIn("match_ratio", validation)
        self.assertIn("all_matched", validation)
        self.assertIn("details", validation)
        self.assertEqual(validation["total"], 4)  # strongest-demo 有 4 个 source units
        # 注意: mock 返回的是 example slide spec（只有 3 个 content slides），
        # 而 normalized 输入是 strongest-demo（4 个 units），所以不一定全部匹配
        self.assertIsInstance(validation["all_matched"], bool)
        self.assertGreater(validation["matched"] + validation["mismatched"], 0)

        # 核心 slide spec 内容（不含 _layout_validation）应与 expected 一致
        drafted_core = {k: v for k, v in drafted.items() if k != "_layout_validation"}
        self.assertEqual(drafted_core, expected)

    def test_openai_compatible_provider_rejects_invalid_slide_spec_shape(self) -> None:
        def fake_transport(_: dict) -> dict:
            return {
                "choices": [
                    {
                        "message": {
                            "content": json.dumps({"deck_goal": "missing slides"}, ensure_ascii=False)
                        }
                    }
                ]
            }

        provider = OpenAICompatibleProvider(
            ProviderConfig(
                provider="openai-compatible",
                model="gpt-4.1-mini",
                api_key_env="GROUNDED_DECK_API_KEY",
                base_url="https://api.example.com/v1",
            ),
            api_key="secret",
            transport=fake_transport,
        )

        with self.assertRaises(ValueError):
            provider.draft_slide_spec(load_json(NORMALIZED_FIXTURE))

    def test_openai_compatible_provider_can_grade_with_mock_transport(self) -> None:
        expected = {
            "status": "pass",
            "failures": [],
            "coverage": {"required_units": 3, "covered_units": 3},
            "grounding": {"total_content_slides": 2, "grounded_slides": 2},
            "visual_form": {"expected_units": 3, "matched_units": 3},
            "provider": "openai-compatible",
            "model": "gpt-4.1-mini",
        }
        normalized = load_json(STRONGEST_DEMO_NORMALIZED_FIXTURE)
        slide_spec = load_json(STRONGEST_DEMO_SLIDE_SPEC_FIXTURE)

        def fake_transport(request: dict) -> dict:
            self.assertIn("Grade strictly against source coverage, slide grounding, and visual-form fit.", request["json"]["messages"][0]["content"])
            self.assertIn("Check whether each unit-backed content slide uses the appropriate layout_type", request["json"]["messages"][1]["content"])
            self.assertIn("Prefer fail over partial credit", request["json"]["messages"][1]["content"])
            self.assertIn("slide_spec", request["json"]["messages"][1]["content"])
            self.assertIn("slide-src-01-sec-01", request["json"]["messages"][1]["content"])
            return {
                "choices": [
                    {
                        "message": {
                            "content": json.dumps(expected, ensure_ascii=False)
                        }
                    }
                ]
            }

        provider = OpenAICompatibleProvider(
            ProviderConfig(
                provider="openai-compatible",
                model="gpt-4.1-mini",
                api_key_env="GROUNDED_DECK_API_KEY",
                base_url="https://api.example.com/v1",
            ),
            api_key="secret",
            transport=fake_transport,
        )

        report = provider.grade_slide_spec(
            normalized,
            slide_spec,
        )
        # 核心 grading 结果（不含 _narrative_validation）应与 expected 一致
        report_core = {k: v for k, v in report.items() if k != "_narrative_validation"}
        self.assertEqual(report_core, expected)

        # 验证 _narrative_validation 附加字段存在且结构正确
        self.assertIn("_narrative_validation", report)
        nv = report["_narrative_validation"]
        self.assertIn("mode", nv)
        self.assertEqual(nv["mode"], "deterministic")
        self.assertIn("slide_count", nv)
        self.assertIn("avg_composite", nv)
        self.assertIn("status", nv)
        self.assertIn("issues", nv)

    def test_openai_compatible_provider_grade_narrative_validation_with_example_fixture(self) -> None:
        """grade_slide_spec 对 example fixture 也能正确附加 _narrative_validation。"""
        expected = {
            "status": "pass",
            "failures": [],
            "coverage": {"required_units": 3, "covered_units": 3},
            "grounding": {"total_content_slides": 2, "grounded_slides": 2},
            "visual_form": {"expected_units": 3, "matched_units": 3},
            "provider": "openai-compatible",
            "model": "gpt-4.1-mini",
        }

        def fake_transport(_: dict) -> dict:
            return {
                "choices": [
                    {"message": {"content": json.dumps(expected, ensure_ascii=False)}}
                ]
            }

        provider = OpenAICompatibleProvider(
            ProviderConfig(
                provider="openai-compatible",
                model="gpt-4.1-mini",
                api_key_env="GROUNDED_DECK_API_KEY",
                base_url="https://api.example.com/v1",
            ),
            api_key="secret",
            transport=fake_transport,
        )

        report = provider.grade_slide_spec(
            load_json(NORMALIZED_FIXTURE),
            load_json(SLIDE_SPEC_FIXTURE),
        )
        self.assertIn("_narrative_validation", report)
        nv = report["_narrative_validation"]
        self.assertEqual(nv["mode"], "deterministic")
        self.assertGreater(nv["slide_count"], 0)
        self.assertGreaterEqual(nv["avg_composite"], 0.0)
        self.assertLessEqual(nv["avg_composite"], 1.0)

    def test_openai_compatible_provider_build_layout_callback(self) -> None:
        """build_layout_callback 返回的 callback 向 LLM 发送单 unit 请求并解析布局字符串。"""
        def fake_transport(request: dict) -> dict:
            # 验证 system prompt 包含布局类型列表
            self.assertIn("timeline", request["json"]["messages"][0]["content"])
            self.assertIn("comparison", request["json"]["messages"][0]["content"])
            # 验证 user prompt 包含 unit 内容
            self.assertIn("出口时间线", request["json"]["messages"][1]["content"])
            return {
                "choices": [
                    {"message": {"content": "  timeline  "}}
                ]
            }

        provider = OpenAICompatibleProvider(
            ProviderConfig(
                provider="openai-compatible",
                model="gpt-4.1-mini",
                api_key_env="GROUNDED_DECK_API_KEY",
                base_url="https://api.example.com/v1",
            ),
            api_key="secret",
            transport=fake_transport,
        )

        callback = provider.build_layout_callback()
        unit = {
            "unit_id": "u1",
            "section_heading": "出口时间线",
            "text": "2022 年到 2024 年的变化",
            "claims": [],
        }
        layout = callback(unit)
        self.assertEqual(layout, "timeline")

    def test_openai_compatible_provider_layout_callback_strips_quotes(self) -> None:
        """build_layout_callback 能正确去除模型返回中的引号包裹。"""
        def fake_transport(request: dict) -> dict:
            return {
                "choices": [
                    {"message": {"content": '"comparison"'}}
                ]
            }

        provider = OpenAICompatibleProvider(
            ProviderConfig(
                provider="openai-compatible",
                model="gpt-4.1-mini",
                api_key_env="GROUNDED_DECK_API_KEY",
                base_url="https://api.example.com/v1",
            ),
            api_key="secret",
            transport=fake_transport,
        )

        callback = provider.build_layout_callback()
        unit = {"section_heading": "对比", "text": "A vs B", "claims": []}
        layout = callback(unit)
        self.assertEqual(layout, "comparison")

    def test_grader_prompt_includes_layout_validation_hint(self) -> None:
        """当有 layout_validation_hint 时，grader user prompt 包含规则引擎分析。"""
        normalized = load_json(STRONGEST_DEMO_NORMALIZED_FIXTURE)
        slide_spec = load_json(STRONGEST_DEMO_SLIDE_SPEC_FIXTURE)

        prompt = OpenAICompatibleProvider.build_grader_user_prompt(
            normalized,
            slide_spec,
            layout_validation_hint="Visual selector rule-based validation: all 4 unit layouts match.",
        )
        self.assertIn("Rule-engine layout analysis", prompt)
        self.assertIn("all 4 unit layouts match", prompt)

    def test_grader_prompt_without_hint_has_no_validation_block(self) -> None:
        """当没有 layout_validation_hint 时，grader prompt 不包含规则引擎分析块。"""
        normalized = load_json(STRONGEST_DEMO_NORMALIZED_FIXTURE)
        slide_spec = load_json(STRONGEST_DEMO_SLIDE_SPEC_FIXTURE)

        prompt = OpenAICompatibleProvider.build_grader_user_prompt(normalized, slide_spec)
        self.assertNotIn("Rule-engine layout analysis", prompt)

    def test_openai_compatible_provider_normalizes_partial_quality_report_shape(self) -> None:
        def fake_transport(_: dict) -> dict:
            return {
                "choices": [
                    {
                        "message": {
                            "content": json.dumps(
                                {
                                    "status": "partial",
                                    "failures": ["mismatch"],
                                    "coverage": {
                                        "required_units": ["u1", "u2"],
                                        "covered_units": ["u1", "u2"],
                                    },
                                    "grounding": {"total_content_slides": 2, "grounded_slides": 2},
                                    "visual_form": {
                                        "expected_units": ["u1", "u2"],
                                        "matched_units": ["u1"],
                                    },
                                    "provider": "grounded-deck",
                                    "model": "quality-grader-v1",
                                },
                                ensure_ascii=False,
                            )
                        }
                    }
                ]
            }

        provider = OpenAICompatibleProvider(
            ProviderConfig(
                provider="openai-compatible",
                model="MiniMax-M2.7",
                api_key_env="GROUNDED_DECK_API_KEY",
                base_url="https://api.minimaxi.com/v1",
            ),
            api_key="secret",
            transport=fake_transport,
        )

        report = provider.grade_slide_spec(
            load_json(NORMALIZED_FIXTURE),
            load_json(SLIDE_SPEC_FIXTURE),
        )

        self.assertEqual(report["status"], "fail")
        self.assertEqual(report["coverage"]["required_units"], 2)
        self.assertEqual(report["coverage"]["covered_units"], 2)
        self.assertEqual(report["coverage"]["required_units_ids"], ["u1", "u2"])
        self.assertEqual(report["visual_form"]["expected_units"], 2)
        self.assertEqual(report["visual_form"]["matched_units"], 1)
        # _narrative_validation 也应存在
        self.assertIn("_narrative_validation", report)

    def test_openai_compatible_provider_rejects_invalid_quality_report_shape(self) -> None:
        def fake_transport(_: dict) -> dict:
            return {
                "choices": [
                    {
                        "message": {
                            "content": json.dumps({"status": "pass", "failures": [], "coverage": {}}, ensure_ascii=False)
                        }
                    }
                ]
            }

        provider = OpenAICompatibleProvider(
            ProviderConfig(
                provider="openai-compatible",
                model="gpt-4.1-mini",
                api_key_env="GROUNDED_DECK_API_KEY",
                base_url="https://api.example.com/v1",
            ),
            api_key="secret",
            transport=fake_transport,
        )

        with self.assertRaises(ValueError):
            provider.grade_slide_spec(
                load_json(NORMALIZED_FIXTURE),
                load_json(SLIDE_SPEC_FIXTURE),
            )

    def test_slide_spec_fixture_has_required_schema_fields(self) -> None:
        schema = load_json(SLIDE_SCHEMA)
        slide_spec = load_json(SLIDE_SPEC_FIXTURE)

        required_root = set(schema["required"])
        self.assertTrue(required_root.issubset(slide_spec.keys()))

        slide_required = set(schema["properties"]["slides"]["items"]["required"])
        for slide in slide_spec["slides"]:
            self.assertTrue(slide_required.issubset(slide.keys()), slide["slide_id"])

    def test_run_pipeline_returns_all_intermediate_outputs(self) -> None:
        result = run_pipeline(load_json(RAW_FIXTURE), provider=DeterministicProvider())

        self.assertEqual(result["normalized_pack"], load_json(NORMALIZED_FIXTURE))
        self.assertEqual(result["slide_spec"], load_json(SLIDE_SPEC_FIXTURE))
        self.assertEqual(result["quality_report"]["status"], "pass")
        self.assertEqual(result["provider"], "deterministic")

    def test_write_pipeline_outputs_writes_expected_files(self) -> None:
        result = run_pipeline(load_json(RAW_FIXTURE), provider=DeterministicProvider())
        with tempfile.TemporaryDirectory() as tmpdir:
            output_dir = Path(tmpdir)
            write_pipeline_outputs(output_dir, result)

            self.assertTrue((output_dir / "normalized-pack.json").exists())
            self.assertTrue((output_dir / "slide-spec.json").exists())
            self.assertTrue((output_dir / "quality-report.json").exists())

    def test_write_verification_summary_writes_metadata_file(self) -> None:
        result = run_pipeline(load_json(RAW_FIXTURE), provider=DeterministicProvider())
        with tempfile.TemporaryDirectory() as tmpdir:
            output_dir = Path(tmpdir)
            write_pipeline_outputs(output_dir, result)
            summary_path = write_verification_summary(
                output_dir=output_dir,
                result=result,
                mode="offline-example",
                input_path=RAW_FIXTURE,
            )

            summary = json.loads(summary_path.read_text(encoding="utf-8"))
            self.assertEqual(summary["mode"], "offline-example")
            self.assertEqual(summary["provider"], "deterministic")
            self.assertEqual(summary["artifacts"]["slide_spec"], str(output_dir / "slide-spec.json"))

    def test_write_strongest_demo_bundle_writes_report_and_artifacts(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            output_dir = Path(tmpdir) / "strongest-demo"
            bundle = write_strongest_demo_bundle(
                input_path=STRONGEST_DEMO_RAW_FIXTURE,
                output_dir=output_dir,
                provider=DeterministicProvider(),
            )

            self.assertTrue((output_dir / "normalized-pack.json").exists())
            self.assertTrue((output_dir / "slide-spec.json").exists())
            self.assertTrue((output_dir / "quality-report.json").exists())
            self.assertTrue((output_dir / "verification-summary.json").exists())
            self.assertTrue((output_dir / "strongest-demo-report.md").exists())
            self.assertTrue((output_dir / "strongest-demo.pptx").exists())
            self.assertEqual(bundle["result"]["quality_report"]["status"], "pass")
            self.assertIn("Success Metrics", bundle["report_path"].read_text(encoding="utf-8"))
            # 验证 PPTX 文件非空
            pptx_size = (output_dir / "strongest-demo.pptx").stat().st_size
            self.assertGreater(pptx_size, 0)


class PipelinePptxIntegrationTests(unittest.TestCase):
    """pipeline render_pptx 参数集成测试。"""

    def test_run_pipeline_with_render_pptx_produces_file(self) -> None:
        """run_pipeline 传入 render_pptx 参数时生成 .pptx 文件。"""
        raw = json.loads(STRONGEST_DEMO_RAW_FIXTURE.read_text(encoding="utf-8"))
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = Path(tmpdir) / "test-output.pptx"
            result = run_pipeline(raw, render_pptx=pptx_path)
            self.assertIn("pptx_path", result)
            self.assertTrue(Path(result["pptx_path"]).exists())
            self.assertTrue(Path(result["pptx_path"]).suffix == ".pptx")

    def test_run_pipeline_without_render_pptx_has_no_pptx_key(self) -> None:
        """run_pipeline 不传 render_pptx 时结果中没有 pptx_path。"""
        raw = json.loads(STRONGEST_DEMO_RAW_FIXTURE.read_text(encoding="utf-8"))
        result = run_pipeline(raw)
        self.assertNotIn("pptx_path", result)

    def test_run_pipeline_render_pptx_produces_valid_pptx(self) -> None:
        """渲染的 PPTX 文件可以被 python-pptx 打开。"""
        from pptx import Presentation

        raw = json.loads(STRONGEST_DEMO_RAW_FIXTURE.read_text(encoding="utf-8"))
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = Path(tmpdir) / "test.pptx"
            result = run_pipeline(raw, render_pptx=pptx_path)
            prs = Presentation(result["pptx_path"])
            self.assertEqual(len(prs.slides), len(result["slide_spec"]["slides"]))

    def test_run_pipeline_render_pptx_creates_directories(self) -> None:
        """render_pptx 路径中的中间目录会自动创建。"""
        raw = json.loads(STRONGEST_DEMO_RAW_FIXTURE.read_text(encoding="utf-8"))
        with tempfile.TemporaryDirectory() as tmpdir:
            deep_path = Path(tmpdir) / "a" / "b" / "c" / "output.pptx"
            result = run_pipeline(raw, render_pptx=deep_path)
            self.assertTrue(Path(result["pptx_path"]).exists())

    def test_run_pipeline_render_pptx_string_path(self) -> None:
        """render_pptx 接受字符串路径。"""
        raw = json.loads(STRONGEST_DEMO_RAW_FIXTURE.read_text(encoding="utf-8"))
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = str(Path(tmpdir) / "test.pptx")
            result = run_pipeline(raw, render_pptx=pptx_path)
            self.assertIn("pptx_path", result)
            self.assertTrue(Path(result["pptx_path"]).exists())

    def test_strongest_demo_bundle_pptx_in_report_text(self) -> None:
        """strongest-demo report 中提及 PPTX 输出。"""
        with tempfile.TemporaryDirectory() as tmpdir:
            output_dir = Path(tmpdir) / "demo"
            bundle = write_strongest_demo_bundle(
                input_path=STRONGEST_DEMO_RAW_FIXTURE,
                output_dir=output_dir,
                provider=DeterministicProvider(),
            )
            report_text = bundle["report_path"].read_text(encoding="utf-8")
            self.assertIn("strongest-demo.pptx", report_text)

    def test_strongest_demo_bundle_pptx_path_in_result(self) -> None:
        """strongest-demo bundle 的 result 中包含 pptx_path。"""
        with tempfile.TemporaryDirectory() as tmpdir:
            output_dir = Path(tmpdir) / "demo"
            bundle = write_strongest_demo_bundle(
                input_path=STRONGEST_DEMO_RAW_FIXTURE,
                output_dir=output_dir,
                provider=DeterministicProvider(),
            )
            self.assertIn("pptx_path", bundle["result"])

    def test_strongest_demo_bundle_pptx_contains_correct_slide_count(self) -> None:
        """strongest-demo 的 PPTX 包含正确的 slide 数量。"""
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as tmpdir:
            output_dir = Path(tmpdir) / "demo"
            bundle = write_strongest_demo_bundle(
                input_path=STRONGEST_DEMO_RAW_FIXTURE,
                output_dir=output_dir,
                provider=DeterministicProvider(),
            )
            pptx_path = output_dir / "strongest-demo.pptx"
            prs = Presentation(str(pptx_path))
            expected_count = len(bundle["result"]["slide_spec"]["slides"])
            self.assertEqual(len(prs.slides), expected_count)


class PipelineArtifactGradingTests(unittest.TestCase):
    """pipeline artifact grading 集成测试。"""

    def test_run_pipeline_with_pptx_includes_artifact_grade(self) -> None:
        """render_pptx 时默认执行 artifact grading。"""
        raw = json.loads(STRONGEST_DEMO_RAW_FIXTURE.read_text(encoding="utf-8"))
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = Path(tmpdir) / "test.pptx"
            result = run_pipeline(raw, render_pptx=pptx_path)
            self.assertIn("artifact_grade", result)
            self.assertEqual(result["artifact_grade"]["status"], "pass")

    def test_run_pipeline_artifact_grade_has_metrics(self) -> None:
        """artifact grade 包含完整的 metrics。"""
        raw = json.loads(STRONGEST_DEMO_RAW_FIXTURE.read_text(encoding="utf-8"))
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = Path(tmpdir) / "test.pptx"
            result = run_pipeline(raw, render_pptx=pptx_path)
            metrics = result["artifact_grade"]["metrics"]
            self.assertIn("slide_count", metrics)
            self.assertIn("editability_ratio", metrics)
            self.assertIn("notes_coverage_ratio", metrics)
            self.assertIn("source_binding_coverage_ratio", metrics)
            self.assertIn("chinese_text_found", metrics)

    def test_run_pipeline_artifact_grade_slide_count_matches_spec(self) -> None:
        """artifact grade 的 slide_count 与 slide spec 一致。"""
        raw = json.loads(STRONGEST_DEMO_RAW_FIXTURE.read_text(encoding="utf-8"))
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = Path(tmpdir) / "test.pptx"
            result = run_pipeline(raw, render_pptx=pptx_path)
            self.assertEqual(
                result["artifact_grade"]["metrics"]["slide_count"],
                len(result["slide_spec"]["slides"]),
            )

    def test_run_pipeline_artifact_grade_disabled(self) -> None:
        """grade_artifact=False 时不执行 artifact grading。"""
        raw = json.loads(STRONGEST_DEMO_RAW_FIXTURE.read_text(encoding="utf-8"))
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = Path(tmpdir) / "test.pptx"
            result = run_pipeline(raw, render_pptx=pptx_path, grade_artifact=False)
            self.assertNotIn("artifact_grade", result)
            self.assertIn("pptx_path", result)

    def test_run_pipeline_no_pptx_no_artifact_grade(self) -> None:
        """不渲染 PPTX 时没有 artifact grade。"""
        raw = json.loads(STRONGEST_DEMO_RAW_FIXTURE.read_text(encoding="utf-8"))
        result = run_pipeline(raw)
        self.assertNotIn("artifact_grade", result)
        self.assertNotIn("pptx_path", result)

    def test_strongest_demo_bundle_includes_artifact_grade(self) -> None:
        """strongest-demo bundle 包含 artifact grade。"""
        with tempfile.TemporaryDirectory() as tmpdir:
            output_dir = Path(tmpdir) / "demo"
            bundle = write_strongest_demo_bundle(
                input_path=STRONGEST_DEMO_RAW_FIXTURE,
                output_dir=output_dir,
                provider=DeterministicProvider(),
            )
            self.assertIn("artifact_grade", bundle["result"])
            self.assertEqual(bundle["result"]["artifact_grade"]["status"], "pass")
            # 验证 artifact-grade.json 文件存在
            self.assertTrue((output_dir / "artifact-grade.json").exists())

    def test_strongest_demo_bundle_report_includes_artifact_section(self) -> None:
        """strongest-demo report 包含 Artifact Grade 部分。"""
        with tempfile.TemporaryDirectory() as tmpdir:
            output_dir = Path(tmpdir) / "demo"
            bundle = write_strongest_demo_bundle(
                input_path=STRONGEST_DEMO_RAW_FIXTURE,
                output_dir=output_dir,
                provider=DeterministicProvider(),
            )
            report_text = bundle["report_path"].read_text(encoding="utf-8")
            self.assertIn("## Artifact Grade", report_text)
            self.assertIn("Editability:", report_text)
            self.assertIn("Notes Coverage:", report_text)

    def test_artifact_grade_editability_is_perfect(self) -> None:
        """所有文本框应可编辑。"""
        raw = json.loads(STRONGEST_DEMO_RAW_FIXTURE.read_text(encoding="utf-8"))
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = Path(tmpdir) / "test.pptx"
            result = run_pipeline(raw, render_pptx=pptx_path)
            self.assertEqual(result["artifact_grade"]["metrics"]["editability_ratio"], 1.0)

    def test_artifact_grade_detects_chinese_text(self) -> None:
        """strongest-demo PPTX 应包含中文。"""
        raw = json.loads(STRONGEST_DEMO_RAW_FIXTURE.read_text(encoding="utf-8"))
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = Path(tmpdir) / "test.pptx"
            result = run_pipeline(raw, render_pptx=pptx_path)
            self.assertTrue(result["artifact_grade"]["metrics"]["chinese_text_found"])


class NarrativePipelineIntegrationTests(unittest.TestCase):
    """pipeline narrative grading 集成测试。"""

    def test_run_pipeline_includes_narrative_grade(self) -> None:
        """pipeline 默认执行 narrative grading。"""
        raw = json.loads(STRONGEST_DEMO_RAW_FIXTURE.read_text(encoding="utf-8"))
        result = run_pipeline(raw)
        self.assertIn("narrative_grade", result)
        self.assertEqual(result["narrative_grade"]["status"], "pass")

    def test_run_pipeline_narrative_grade_has_required_fields(self) -> None:
        """narrative grade 包含完整的字段。"""
        raw = json.loads(STRONGEST_DEMO_RAW_FIXTURE.read_text(encoding="utf-8"))
        result = run_pipeline(raw)
        ng = result["narrative_grade"]
        self.assertIn("status", ng)
        self.assertIn("mode", ng)
        self.assertIn("avg_coherence", ng)
        self.assertIn("avg_grounding", ng)
        self.assertIn("avg_visual_fit", ng)
        self.assertIn("avg_composite", ng)
        self.assertIn("slides", ng)

    def test_run_pipeline_narrative_grade_slide_count_matches(self) -> None:
        """narrative grade 的 slide_count 与 slide spec 一致。"""
        raw = json.loads(STRONGEST_DEMO_RAW_FIXTURE.read_text(encoding="utf-8"))
        result = run_pipeline(raw)
        self.assertEqual(
            result["narrative_grade"]["slide_count"],
            len(result["slide_spec"]["slides"]),
        )

    def test_run_pipeline_narrative_grade_disabled(self) -> None:
        """grade_narrative_quality=False 时不执行 narrative grading。"""
        raw = json.loads(STRONGEST_DEMO_RAW_FIXTURE.read_text(encoding="utf-8"))
        result = run_pipeline(raw, grade_narrative_quality=False)
        self.assertNotIn("narrative_grade", result)

    def test_run_pipeline_narrative_grade_avg_composite_above_threshold(self) -> None:
        """strongest-demo 的 avg_composite 应 >= 0.7。"""
        raw = json.loads(STRONGEST_DEMO_RAW_FIXTURE.read_text(encoding="utf-8"))
        result = run_pipeline(raw)
        self.assertGreaterEqual(result["narrative_grade"]["avg_composite"], 0.7)

    def test_strongest_demo_bundle_includes_narrative_grade(self) -> None:
        """strongest-demo bundle 包含 narrative grade。"""
        with tempfile.TemporaryDirectory() as tmpdir:
            output_dir = Path(tmpdir) / "demo"
            bundle = write_strongest_demo_bundle(
                input_path=STRONGEST_DEMO_RAW_FIXTURE,
                output_dir=output_dir,
                provider=DeterministicProvider(),
            )
            self.assertIn("narrative_grade", bundle["result"])
            self.assertEqual(bundle["result"]["narrative_grade"]["status"], "pass")
            # 验证 narrative-grade.json 文件存在
            self.assertTrue((output_dir / "narrative-grade.json").exists())

    def test_strongest_demo_bundle_report_includes_narrative_section(self) -> None:
        """strongest-demo report 包含 Narrative Grade 部分。"""
        with tempfile.TemporaryDirectory() as tmpdir:
            output_dir = Path(tmpdir) / "demo"
            bundle = write_strongest_demo_bundle(
                input_path=STRONGEST_DEMO_RAW_FIXTURE,
                output_dir=output_dir,
                provider=DeterministicProvider(),
            )
            report_text = bundle["report_path"].read_text(encoding="utf-8")
            self.assertIn("## Narrative Grade", report_text)
            self.assertIn("Avg Coherence:", report_text)
            self.assertIn("Avg Grounding:", report_text)
            self.assertIn("Avg Visual Fit:", report_text)


if __name__ == "__main__":
    unittest.main()
