"""GroundedDeck PPTX 渲染器测试。

覆盖：
- 基本渲染流程
- 所有 7 种布局类型
- strongest-demo fixture 的完整回归
- 错误处理
- 输出文件的可编辑性验证
"""

from __future__ import annotations

import json
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from src.renderer.pptx_renderer import (
    get_supported_layouts,
    render_slide_spec_to_pptx,
)

ROOT = Path(__file__).resolve().parent.parent
STRONGEST_DEMO_SLIDE_SPEC = ROOT / "fixtures" / "slide-spec" / "strongest-demo-slide-spec.json"
STRONGEST_DEMO_SOURCE_PACK = ROOT / "fixtures" / "source-packs" / "strongest-demo-source-pack.json"
EXAMPLE_SLIDE_SPEC = ROOT / "fixtures" / "slide-spec" / "example-slide-spec.json"


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def _make_minimal_spec(slides: list[dict]) -> dict:
    """构造一个最小的 slide spec。"""
    return {
        "deck_goal": "Test goal",
        "audience": "Test audience",
        "slides": slides,
    }


def _make_slide(
    slide_id: str = "s1",
    title: str = "Test Title",
    layout_type: str = "summary",
    **kwargs,
) -> dict:
    """构造一个最小的 slide dict。"""
    slide = {
        "slide_id": slide_id,
        "title": title,
        "goal": kwargs.get("goal", "Test goal"),
        "layout_type": layout_type,
        "key_points": kwargs.get("key_points", ["Point 1"]),
        "visual_elements": kwargs.get("visual_elements", [{"type": "bullet-list"}]),
        "source_bindings": kwargs.get("source_bindings", ["src-01:sec-01"]),
        "must_include_checks": kwargs.get("must_include_checks", ["src-01:sec-01"]),
    }
    if "speaker_notes" in kwargs:
        slide["speaker_notes"] = kwargs["speaker_notes"]
    return slide


class BasicRenderingTests(unittest.TestCase):
    """基本渲染流程测试。"""

    def test_render_minimal_spec_produces_pptx_file(self) -> None:
        """最小的 slide spec 能产生 .pptx 文件。"""
        spec = _make_minimal_spec([_make_slide()])
        with tempfile.TemporaryDirectory() as tmpdir:
            out = render_slide_spec_to_pptx(spec, Path(tmpdir) / "test.pptx")
            self.assertTrue(out.exists())
            self.assertTrue(out.suffix == ".pptx")

    def test_render_produces_correct_slide_count(self) -> None:
        """渲染器生成的 slide 数量与 spec 一致。"""
        slides = [_make_slide(slide_id=f"s{i}") for i in range(3)]
        spec = _make_minimal_spec(slides)
        with tempfile.TemporaryDirectory() as tmpdir:
            out = render_slide_spec_to_pptx(spec, Path(tmpdir) / "test.pptx")
            prs = Presentation(str(out))
            self.assertEqual(len(prs.slides), 3)

    def test_render_creates_output_directory_if_missing(self) -> None:
        """输出路径不存在时自动创建。"""
        spec = _make_minimal_spec([_make_slide()])
        with tempfile.TemporaryDirectory() as tmpdir:
            deep_path = Path(tmpdir) / "a" / "b" / "c" / "test.pptx"
            out = render_slide_spec_to_pptx(spec, deep_path)
            self.assertTrue(out.exists())

    def test_render_returns_path_object(self) -> None:
        """返回值是 Path 对象。"""
        spec = _make_minimal_spec([_make_slide()])
        with tempfile.TemporaryDirectory() as tmpdir:
            out = render_slide_spec_to_pptx(spec, Path(tmpdir) / "test.pptx")
            self.assertIsInstance(out, Path)

    def test_empty_slides_raises_value_error(self) -> None:
        """空 slides 列表抛出 ValueError。"""
        spec = _make_minimal_spec([])
        with tempfile.TemporaryDirectory() as tmpdir:
            with self.assertRaises(ValueError):
                render_slide_spec_to_pptx(spec, Path(tmpdir) / "test.pptx")

    def test_missing_slides_key_raises_value_error(self) -> None:
        """缺少 slides 键抛出 ValueError。"""
        spec = {"deck_goal": "Test", "audience": "Test"}
        with tempfile.TemporaryDirectory() as tmpdir:
            with self.assertRaises(ValueError):
                render_slide_spec_to_pptx(spec, Path(tmpdir) / "test.pptx")

    def test_speaker_notes_written_to_notes_slide(self) -> None:
        """speaker notes 写入 slide 的备注区。"""
        spec = _make_minimal_spec([
            _make_slide(speaker_notes="This is a test note."),
        ])
        with tempfile.TemporaryDirectory() as tmpdir:
            out = render_slide_spec_to_pptx(spec, Path(tmpdir) / "test.pptx")
            prs = Presentation(str(out))
            notes = prs.slides[0].notes_slide.notes_text_frame.text
            self.assertIn("This is a test note.", notes)

    def test_source_bindings_in_notes(self) -> None:
        """source bindings 出现在备注中。"""
        spec = _make_minimal_spec([
            _make_slide(source_bindings=["src-01:sec-01", "src-02:sec-01"]),
        ])
        with tempfile.TemporaryDirectory() as tmpdir:
            out = render_slide_spec_to_pptx(spec, Path(tmpdir) / "test.pptx")
            prs = Presentation(str(out))
            notes = prs.slides[0].notes_slide.notes_text_frame.text
            self.assertIn("[Source Bindings]", notes)
            self.assertIn("src-01:sec-01", notes)
            self.assertIn("src-02:sec-01", notes)


class LayoutRenderingTests(unittest.TestCase):
    """每种布局类型的渲染测试。"""

    def _render_single(self, layout_type: str, **kwargs) -> Presentation:
        """渲染单个指定布局的 slide 并返回 Presentation 对象。"""
        spec = _make_minimal_spec([_make_slide(layout_type=layout_type, **kwargs)])
        with tempfile.TemporaryDirectory() as tmpdir:
            out = render_slide_spec_to_pptx(spec, Path(tmpdir) / "test.pptx")
            return Presentation(str(out))

    def test_cover_layout_renders(self) -> None:
        """cover 布局能正常渲染。"""
        prs = self._render_single(
            "cover",
            title="Test Cover",
            key_points=["Audience info", "4 sources"],
            source_bindings=["src-01:sec-01"],
        )
        self.assertEqual(len(prs.slides), 1)
        # cover 页应该有多个 shapes
        self.assertGreater(len(prs.slides[0].shapes), 0)

    def test_summary_layout_renders(self) -> None:
        """summary 布局能正常渲染。"""
        prs = self._render_single(
            "summary",
            key_points=["Point 1", "Point 2", "Point 3"],
        )
        self.assertEqual(len(prs.slides), 1)
        self.assertGreater(len(prs.slides[0].shapes), 0)

    def test_timeline_layout_renders(self) -> None:
        """timeline 布局能正常渲染，包含里程碑。"""
        prs = self._render_single(
            "timeline",
            visual_elements=[{"type": "timeline", "milestones": ["2022", "2023", "2024"]}],
            key_points=["出口节奏变化"],
        )
        self.assertEqual(len(prs.slides), 1)
        # timeline 至少有：标题 + 时间线 + 3个节点 + 3个标签 + key_point = 多个 shapes
        self.assertGreater(len(prs.slides[0].shapes), 5)

    def test_comparison_layout_renders(self) -> None:
        """comparison 布局能正常渲染，包含双栏。"""
        prs = self._render_single(
            "comparison",
            visual_elements=[{"type": "comparison-columns", "columns": ["欧洲", "东南亚"]}],
            key_points=["对比分析"],
        )
        self.assertEqual(len(prs.slides), 1)
        self.assertGreater(len(prs.slides[0].shapes), 3)

    def test_process_layout_renders(self) -> None:
        """process 布局能正常渲染。"""
        prs = self._render_single(
            "process",
            visual_elements=[{"type": "process-flow", "steps": 1}],
            key_points=["先建立渠道再进入欧洲"],
        )
        self.assertEqual(len(prs.slides), 1)
        self.assertGreater(len(prs.slides[0].shapes), 0)

    def test_chart_layout_renders(self) -> None:
        """chart 布局能正常渲染，包含指标卡片。"""
        prs = self._render_single(
            "chart",
            visual_elements=[{"type": "metric-cards", "metrics": ["12", "18%"]}],
            key_points=["成本降幅"],
        )
        self.assertEqual(len(prs.slides), 1)
        # chart: 标题 + 2个卡片背景 + 2个数字 + 2个标签 + key_point
        self.assertGreater(len(prs.slides[0].shapes), 5)

    def test_section_layout_renders(self) -> None:
        """section 布局能正常渲染。"""
        prs = self._render_single("section", goal="New section begins")
        self.assertEqual(len(prs.slides), 1)
        self.assertGreater(len(prs.slides[0].shapes), 0)

    def test_unsupported_layout_falls_back_to_summary(self) -> None:
        """不支持的 layout_type 回退到 summary。"""
        prs = self._render_single("fancy-3d-view")
        self.assertEqual(len(prs.slides), 1)
        # 应该能渲染成功（作为 summary）
        self.assertGreater(len(prs.slides[0].shapes), 0)

    def test_timeline_without_milestones_uses_default(self) -> None:
        """timeline 布局没有 milestones 时使用默认值。"""
        prs = self._render_single(
            "timeline",
            visual_elements=[{"type": "bullet-list"}],  # 没有 timeline 类型的 VE
        )
        self.assertEqual(len(prs.slides), 1)

    def test_comparison_with_less_than_two_columns_uses_default(self) -> None:
        """comparison 布局列数不足时使用默认值。"""
        prs = self._render_single(
            "comparison",
            visual_elements=[{"type": "comparison-columns", "columns": ["OnlyOne"]}],
        )
        self.assertEqual(len(prs.slides), 1)

    def test_chart_without_metrics_uses_default(self) -> None:
        """chart 布局没有 metrics 时使用默认值。"""
        prs = self._render_single(
            "chart",
            visual_elements=[{"type": "bullet-list"}],
        )
        self.assertEqual(len(prs.slides), 1)


class StrongestDemoRegressionTests(unittest.TestCase):
    """strongest-demo fixture 的完整回归测试。"""

    def test_strongest_demo_renders_all_6_slides(self) -> None:
        """strongest-demo slide spec 渲染出正确数量的 slides。"""
        spec = load_json(STRONGEST_DEMO_SLIDE_SPEC)
        with tempfile.TemporaryDirectory() as tmpdir:
            out = render_slide_spec_to_pptx(spec, Path(tmpdir) / "strongest-demo.pptx")
            prs = Presentation(str(out))
            self.assertEqual(len(prs.slides), 6)

    def test_strongest_demo_layout_sequence(self) -> None:
        """strongest-demo 的布局序列正确。"""
        spec = load_json(STRONGEST_DEMO_SLIDE_SPEC)
        expected_layouts = ["cover", "summary", "timeline", "comparison", "process", "chart"]
        actual_layouts = [s["layout_type"] for s in spec["slides"]]
        self.assertEqual(actual_layouts, expected_layouts)

    def test_strongest_demo_output_is_valid_pptx(self) -> None:
        """strongest-demo 输出文件是有效的 .pptx，可以被 python-pptx 重新打开。"""
        spec = load_json(STRONGEST_DEMO_SLIDE_SPEC)
        with tempfile.TemporaryDirectory() as tmpdir:
            out = render_slide_spec_to_pptx(spec, Path(tmpdir) / "strongest-demo.pptx")
            # 验证可以重新打开
            prs = Presentation(str(out))
            self.assertEqual(len(prs.slides), 6)
            # 验证宽屏比例
            self.assertGreater(prs.slide_width, prs.slide_height)

    def test_strongest_demo_all_slides_have_shapes(self) -> None:
        """strongest-demo 每个 slide 都有渲染内容。"""
        spec = load_json(STRONGEST_DEMO_SLIDE_SPEC)
        with tempfile.TemporaryDirectory() as tmpdir:
            out = render_slide_spec_to_pptx(spec, Path(tmpdir) / "strongest-demo.pptx")
            prs = Presentation(str(out))
            for i, slide in enumerate(prs.slides):
                self.assertGreater(
                    len(slide.shapes),
                    0,
                    f"slide {i} has no shapes",
                )

    def test_strongest_demo_all_slides_have_notes(self) -> None:
        """strongest-demo 每个 slide 都有备注。"""
        spec = load_json(STRONGEST_DEMO_SLIDE_SPEC)
        with tempfile.TemporaryDirectory() as tmpdir:
            out = render_slide_spec_to_pptx(spec, Path(tmpdir) / "strongest-demo.pptx")
            prs = Presentation(str(out))
            for i, slide in enumerate(prs.slides):
                notes_text = slide.notes_slide.notes_text_frame.text
                self.assertTrue(
                    len(notes_text) > 0,
                    f"slide {i} has empty notes",
                )

    def test_strongest_demo_cover_slide_has_source_count(self) -> None:
        """strongest-demo 封面页有 source count 信息。"""
        spec = load_json(STRONGEST_DEMO_SLIDE_SPEC)
        cover = spec["slides"][0]
        self.assertEqual(cover["layout_type"], "cover")
        self.assertEqual(len(cover["source_bindings"]), 4)

    def test_strongest_demo_file_size_reasonable(self) -> None:
        """输出文件大小在合理范围内（不应该超过 5MB 或为 0）。"""
        spec = load_json(STRONGEST_DEMO_SLIDE_SPEC)
        with tempfile.TemporaryDirectory() as tmpdir:
            out = render_slide_spec_to_pptx(spec, Path(tmpdir) / "strongest-demo.pptx")
            file_size = out.stat().st_size
            self.assertGreater(file_size, 0)
            self.assertLess(file_size, 5 * 1024 * 1024)  # < 5 MB


class SupportedLayoutsTests(unittest.TestCase):
    """get_supported_layouts API 测试。"""

    def test_returns_tuple_of_strings(self) -> None:
        layouts = get_supported_layouts()
        self.assertIsInstance(layouts, tuple)
        for layout in layouts:
            self.assertIsInstance(layout, str)

    def test_contains_all_expected_layouts(self) -> None:
        layouts = get_supported_layouts()
        expected = {"cover", "summary", "timeline", "comparison", "process", "chart", "section"}
        self.assertEqual(set(layouts), expected)


class ExampleFixtureTests(unittest.TestCase):
    """example slide spec fixture 渲染测试。"""

    def test_example_spec_renders_without_error(self) -> None:
        """example slide spec 也能成功渲染。"""
        spec = load_json(EXAMPLE_SLIDE_SPEC)
        with tempfile.TemporaryDirectory() as tmpdir:
            out = render_slide_spec_to_pptx(spec, Path(tmpdir) / "example.pptx")
            prs = Presentation(str(out))
            self.assertEqual(len(prs.slides), len(spec["slides"]))


# ================================================================
# Artifact Grading Tests
# ================================================================


class ArtifactGradingTests(unittest.TestCase):
    """PPTX artifact quality grading 测试。"""

    def _render_spec(self, spec: dict) -> Path:
        """辅助方法：渲染 spec 到临时文件。"""
        self._tmpdir = tempfile.mkdtemp()
        out = render_slide_spec_to_pptx(spec, Path(self._tmpdir) / "test.pptx")
        return out

    def test_grade_strongest_demo_passes(self) -> None:
        """strongest demo 渲染结果应通过 artifact grading。"""
        from src.renderer.artifact_grader import grade_pptx_artifact

        spec = load_json(STRONGEST_DEMO_SLIDE_SPEC)
        out = self._render_spec(spec)
        report = grade_pptx_artifact(out, slide_spec=spec)
        self.assertEqual(report["status"], "pass", f"failures: {report['failures']}")
        self.assertEqual(report["metrics"]["slide_count"], len(spec["slides"]))
        self.assertTrue(report["metrics"]["chinese_text_found"])

    def test_grade_artifact_detects_slide_count_mismatch(self) -> None:
        """当 slide count 不匹配 spec 时应失败。"""
        from src.renderer.artifact_grader import grade_pptx_artifact

        spec = load_json(STRONGEST_DEMO_SLIDE_SPEC)
        out = self._render_spec(spec)
        # 构造一个少一个 slide 的假 spec
        fake_spec = dict(spec)
        fake_spec["slides"] = spec["slides"][:3]
        report = grade_pptx_artifact(out, slide_spec=fake_spec)
        self.assertEqual(report["status"], "fail")
        self.assertTrue(any("slide count mismatch" in f for f in report["failures"]))

    def test_grade_artifact_missing_file(self) -> None:
        """不存在的文件应失败。"""
        from src.renderer.artifact_grader import grade_pptx_artifact

        report = grade_pptx_artifact("/nonexistent/path.pptx")
        self.assertEqual(report["status"], "fail")
        self.assertTrue(any("not found" in f for f in report["failures"]))

    def test_grade_artifact_editability(self) -> None:
        """所有文本框应可编辑。"""
        from src.renderer.artifact_grader import grade_pptx_artifact

        spec = load_json(STRONGEST_DEMO_SLIDE_SPEC)
        out = self._render_spec(spec)
        report = grade_pptx_artifact(out)
        self.assertEqual(report["metrics"]["editability_ratio"], 1.0)
        self.assertGreater(report["metrics"]["editable_text_boxes"], 0)

    def test_grade_artifact_notes_coverage(self) -> None:
        """所有 slides 应有 speaker notes。"""
        from src.renderer.artifact_grader import grade_pptx_artifact

        spec = load_json(STRONGEST_DEMO_SLIDE_SPEC)
        out = self._render_spec(spec)
        report = grade_pptx_artifact(out)
        self.assertEqual(
            report["metrics"]["slides_with_notes"],
            report["metrics"]["slide_count"],
        )
        self.assertEqual(report["metrics"]["notes_coverage_ratio"], 1.0)

    def test_grade_artifact_source_bindings_in_notes(self) -> None:
        """大多数 slides 应在备注中包含 source bindings。"""
        from src.renderer.artifact_grader import grade_pptx_artifact

        spec = load_json(STRONGEST_DEMO_SLIDE_SPEC)
        out = self._render_spec(spec)
        report = grade_pptx_artifact(out)
        self.assertGreater(report["metrics"]["slides_with_source_bindings"], 0)
        self.assertGreaterEqual(report["metrics"]["source_binding_coverage_ratio"], 0.5)


# ================================================================
# Narrative Quality Grading Tests
# ================================================================


class NarrativeQualityTests(unittest.TestCase):
    """Narrative quality grading 维度测试。"""

    def test_strongest_demo_narrative_quality_passes(self) -> None:
        """strongest demo 的 narrative quality 应全部通过。"""
        from src.ingest.normalize import normalize_source_pack
        from src.llm.provider import DeterministicProvider

        raw = load_json(STRONGEST_DEMO_SOURCE_PACK)
        norm = normalize_source_pack(raw)
        provider = DeterministicProvider()
        spec = provider.draft_slide_spec(norm)
        report = provider.grade_slide_spec(norm, spec)

        nq = report["narrative_quality"]
        self.assertEqual(nq["slides_with_empty_key_points"], [])
        self.assertEqual(nq["issues"], [])
        self.assertEqual(nq["quality_ratio"], 1.0)

    def test_narrative_quality_detects_empty_key_points(self) -> None:
        """缺少 key_points 的 content slide 应被检出。"""
        from src.ingest.normalize import normalize_source_pack
        from src.llm.provider import DeterministicProvider

        raw = load_json(STRONGEST_DEMO_SOURCE_PACK)
        norm = normalize_source_pack(raw)
        provider = DeterministicProvider()
        spec = provider.draft_slide_spec(norm)

        # 手动清空一个 content slide 的 key_points
        spec["slides"][2]["key_points"] = []
        report = provider.grade_slide_spec(norm, spec)

        nq = report["narrative_quality"]
        self.assertIn(spec["slides"][2]["slide_id"], nq["slides_with_empty_key_points"])
        self.assertGreater(len(report["failures"]), 0)

    def test_narrative_quality_source_annotations_count(self) -> None:
        """summary slide 的 key_points 应包含 source annotations。"""
        from src.ingest.normalize import normalize_source_pack
        from src.llm.provider import DeterministicProvider

        raw = load_json(STRONGEST_DEMO_SOURCE_PACK)
        norm = normalize_source_pack(raw)
        provider = DeterministicProvider()
        spec = provider.draft_slide_spec(norm)
        report = provider.grade_slide_spec(norm, spec)

        nq = report["narrative_quality"]
        # 4 个 source units 的 claims 带有 [binding] 标注
        self.assertEqual(nq["source_annotated_points"], 4)

    def test_cover_key_points_enriched(self) -> None:
        """cover slide 应包含 source-grounded 核心声明。"""
        from src.ingest.normalize import normalize_source_pack
        from src.llm.provider import DeterministicProvider

        raw = load_json(STRONGEST_DEMO_SOURCE_PACK)
        norm = normalize_source_pack(raw)
        provider = DeterministicProvider()
        spec = provider.draft_slide_spec(norm)

        cover = spec["slides"][0]
        self.assertEqual(cover["layout_type"], "cover")
        # cover 应有 3 个 key_points：audience, 核心声明, source count
        self.assertEqual(len(cover["key_points"]), 3)
        # 应包含来自 source units 的内容
        self.assertTrue(any("grounded source units" in kp for kp in cover["key_points"]))

    def test_summary_key_points_cover_all_units(self) -> None:
        """summary slide 的 key_points 应覆盖所有 source units。"""
        from src.ingest.normalize import normalize_source_pack
        from src.llm.provider import DeterministicProvider

        raw = load_json(STRONGEST_DEMO_SOURCE_PACK)
        norm = normalize_source_pack(raw)
        provider = DeterministicProvider()
        spec = provider.draft_slide_spec(norm)

        summary = spec["slides"][1]
        self.assertEqual(summary["layout_type"], "summary")
        # summary 应有 4 个 key_points（每个 unit 1 条 claim）
        self.assertEqual(len(summary["key_points"]), 4)
        # 每条 key_point 应包含 source binding 标注
        for kp in summary["key_points"]:
            self.assertIn("[", kp)
            self.assertIn("]", kp)

    def test_summary_visual_elements_enriched(self) -> None:
        """summary slide 应包含 claim-source-map visual element。"""
        from src.ingest.normalize import normalize_source_pack
        from src.llm.provider import DeterministicProvider

        raw = load_json(STRONGEST_DEMO_SOURCE_PACK)
        norm = normalize_source_pack(raw)
        provider = DeterministicProvider()
        spec = provider.draft_slide_spec(norm)

        summary = spec["slides"][1]
        ve_types = [ve["type"] for ve in summary["visual_elements"]]
        self.assertIn("bullet-list", ve_types)
        self.assertIn("claim-source-map", ve_types)

        # claim-source-map 应有 4 个 entries
        csm = next(ve for ve in summary["visual_elements"] if ve["type"] == "claim-source-map")
        self.assertEqual(len(csm["entries"]), 4)

    def test_cover_visual_elements_enriched(self) -> None:
        """cover slide 应包含 topic-overview visual element。"""
        from src.ingest.normalize import normalize_source_pack
        from src.llm.provider import DeterministicProvider

        raw = load_json(STRONGEST_DEMO_SOURCE_PACK)
        norm = normalize_source_pack(raw)
        provider = DeterministicProvider()
        spec = provider.draft_slide_spec(norm)

        cover = spec["slides"][0]
        ve_types = [ve["type"] for ve in cover["visual_elements"]]
        self.assertIn("topic-overview", ve_types)

        # topic-overview 应包含 4 个 topics
        topics_ve = next(ve for ve in cover["visual_elements"] if ve["type"] == "topic-overview")
        self.assertEqual(len(topics_ve["topics"]), 4)


if __name__ == "__main__":
    unittest.main()
