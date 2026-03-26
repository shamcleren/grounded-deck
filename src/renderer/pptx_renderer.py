"""GroundedDeck PPTX 渲染器：将 slide spec 转换为可编辑的 .pptx 文件。

该模块消费 slide spec JSON 并使用 python-pptx 生成结构化的 PowerPoint 文件。
每种 layout_type 都有独立的渲染函数，确保输出可审计且可编辑。

设计原则：
- 零外部模板依赖，纯代码生成
- 每页 slide 保留 source bindings 和 speaker notes 以支持审计
- 所有文本框保持可编辑状态
- 专业蓝灰色调配色方案
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

import platform

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Emu, Pt

logger = logging.getLogger(__name__)

# ---------- 配色方案 ----------

COLOR_PRIMARY = RGBColor(0x1A, 0x56, 0xDB)      # 主色：深蓝
COLOR_ACCENT = RGBColor(0x2D, 0x8C, 0xF0)        # 辅助色：亮蓝
COLOR_DARK = RGBColor(0x1E, 0x29, 0x3B)           # 深色文字
COLOR_BODY = RGBColor(0x37, 0x41, 0x51)            # 正文文字
COLOR_MUTED = RGBColor(0x6B, 0x72, 0x80)           # 次要文字
COLOR_LIGHT_BG = RGBColor(0xF3, 0xF4, 0xF6)       # 浅色背景
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)           # 白色
COLOR_BORDER = RGBColor(0xD1, 0xD5, 0xDB)          # 边框色
COLOR_HIGHLIGHT = RGBColor(0xDC, 0x26, 0x26)       # 高亮/强调色

# 布局专用颜色
COLOR_TIMELINE_ACCENT = RGBColor(0x05, 0x96, 0x69)   # 时间线：绿色
COLOR_COMPARISON_LEFT = RGBColor(0x1A, 0x56, 0xDB)    # 对比左：蓝色
COLOR_COMPARISON_RIGHT = RGBColor(0xDC, 0x26, 0x26)   # 对比右：红色
COLOR_PROCESS_ACCENT = RGBColor(0x7C, 0x3A, 0xED)     # 流程：紫色
COLOR_CHART_ACCENT = RGBColor(0xD9, 0x77, 0x06)       # 图表：橙色

# ---------- 字体配置（跨平台中文安全回退链） ----------


def _detect_cjk_font() -> tuple[str, str]:
    """根据操作系统检测可用的中文字体，返回 (标题字体, 正文字体)。

    回退优先级：
    - macOS: PingFang SC > Hiragino Sans GB > STHeiti > Microsoft YaHei
    - Windows: Microsoft YaHei > SimHei > DengXian
    - Linux: Noto Sans CJK SC > WenQuanYi Micro Hei > Microsoft YaHei
    """
    system = platform.system()
    if system == "Darwin":  # macOS
        return ("PingFang SC", "PingFang SC")
    elif system == "Windows":
        return ("Microsoft YaHei", "Microsoft YaHei")
    else:  # Linux 及其他
        return ("Noto Sans CJK SC", "Noto Sans CJK SC")


# 主字体 + 回退列表（用于 XML 级别的字体回退）
_CJK_TITLE, _CJK_BODY = _detect_cjk_font()

# 跨平台回退链：渲染时写入主字体，PPTX 打开时由 PowerPoint 自动回退
CJK_FONT_FALLBACK_CHAIN = [
    "PingFang SC",        # macOS
    "Microsoft YaHei",    # Windows
    "Noto Sans CJK SC",   # Linux
    "SimHei",             # Windows 备选
    "Hiragino Sans GB",   # macOS 备选
    "WenQuanYi Micro Hei",  # Linux 备选
]

FONT_TITLE = _CJK_TITLE   # 中文标题字体（平台自适应）
FONT_BODY = _CJK_BODY      # 中文正文字体（平台自适应）
FONT_MONO = "Consolas"     # 等宽字体

# 西文回退字体（当内容为纯英文时使用）
FONT_LATIN_TITLE = "Calibri"
FONT_LATIN_BODY = "Calibri"

# ---------- 页面尺寸（宽屏 16:9） ----------

SLIDE_WIDTH = Cm(33.867)   # 13.333 inches
SLIDE_HEIGHT = Cm(19.05)   # 7.5 inches

# ---------- 常用边距和尺寸 ----------

MARGIN_LEFT = Cm(2.0)
MARGIN_TOP = Cm(2.0)
MARGIN_RIGHT = Cm(2.0)
CONTENT_WIDTH = Cm(29.867)  # SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
TITLE_HEIGHT = Cm(2.5)
CONTENT_TOP = Cm(5.0)       # 标题区域下方


# ================================================================
# 辅助函数
# ================================================================


def _set_font(run, *, size: int = 14, bold: bool = False, color: RGBColor = COLOR_BODY, name: str = FONT_BODY) -> None:
    """设置 run 的字体属性，同时配置东亚字体回退。"""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    # 设置东亚字体（确保中文字符正确渲染）
    _set_east_asian_font(run, name)


def _set_east_asian_font(run, font_name: str) -> None:
    """通过 XML 操作设置 run 的东亚字体属性。

    python-pptx 不直接支持 ea (East Asian) 字体设置，
    需要通过底层 XML 操作来确保中文字符使用正确的字体。
    """
    from lxml import etree
    rPr = run._r.get_or_add_rPr()
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    # 移除已有的 ea 元素
    for ea in rPr.findall("a:ea", nsmap):
        rPr.remove(ea)
    # 添加东亚字体
    ea = etree.SubElement(
        rPr,
        "{http://schemas.openxmlformats.org/drawingml/2006/main}ea",
    )
    ea.set("typeface", font_name)


def _add_textbox(
    slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    text: str,
    *,
    font_size: int = 14,
    bold: bool = False,
    color: RGBColor = COLOR_BODY,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str = FONT_BODY,
    word_wrap: bool = True,
) -> Any:
    """添加一个简单文本框。"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _set_font(run, size=font_size, bold=bold, color=color, name=font_name)
    return txbox


def _add_title_block(slide, title: str, subtitle: str = "") -> None:
    """在 slide 顶部添加标题区域。"""
    # 标题
    _add_textbox(
        slide,
        MARGIN_LEFT,
        MARGIN_TOP,
        CONTENT_WIDTH,
        Cm(1.5),
        title,
        font_size=28,
        bold=True,
        color=COLOR_DARK,
        font_name=FONT_TITLE,
    )
    # 副标题
    if subtitle:
        _add_textbox(
            slide,
            MARGIN_LEFT,
            Cm(3.8),
            CONTENT_WIDTH,
            Cm(1.0),
            subtitle,
            font_size=14,
            color=COLOR_MUTED,
        )


def _add_speaker_notes(slide, slide_data: dict) -> None:
    """将 speaker notes 和 source bindings 写入备注区。"""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame

    # speaker notes
    speaker_notes = slide_data.get("speaker_notes", "")
    if speaker_notes:
        tf.text = speaker_notes

    # source bindings（追加到备注末尾）
    bindings = slide_data.get("source_bindings", [])
    if bindings:
        p = tf.add_paragraph()
        p.text = ""
        p = tf.add_paragraph()
        p.text = f"[Source Bindings] {', '.join(bindings)}"

    # must_include_checks（追加到备注末尾）
    checks = slide_data.get("must_include_checks", [])
    if checks:
        p = tf.add_paragraph()
        p.text = f"[Must Include] {', '.join(checks)}"


def _add_rounded_rect(
    slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    fill_color: RGBColor = COLOR_LIGHT_BG,
    line_color: RGBColor | None = None,
) -> Any:
    """添加一个圆角矩形作为背景卡片。"""
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _add_bullet_list(slide, left: Emu, top: Emu, width: Emu, height: Emu, items: list[str]) -> Any:
    """添加一个带要点符号的文本框。"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"• {item}"
        _set_font(run, size=16, color=COLOR_BODY)

    return txbox


# ================================================================
# 布局渲染器
# ================================================================


def _render_cover(slide, slide_data: dict) -> None:
    """渲染 cover 布局：大标题 + 目标受众 + source count。"""
    # 背景色块
    _add_rounded_rect(slide, Cm(0), Cm(0), SLIDE_WIDTH, Cm(10), fill_color=COLOR_PRIMARY)

    # 主标题（白色大字）
    _add_textbox(
        slide,
        Cm(3.0),
        Cm(3.0),
        Cm(27.0),
        Cm(4.0),
        slide_data["title"],
        font_size=36,
        bold=True,
        color=COLOR_WHITE,
        font_name=FONT_TITLE,
        alignment=PP_ALIGN.LEFT,
    )

    # key_points 作为副信息
    key_points = slide_data.get("key_points", [])
    if key_points:
        subtitle_text = " | ".join(key_points)
        _add_textbox(
            slide,
            Cm(3.0),
            Cm(7.5),
            Cm(27.0),
            Cm(1.5),
            subtitle_text,
            font_size=16,
            color=RGBColor(0xBF, 0xDB, 0xFE),  # 浅蓝色
            alignment=PP_ALIGN.LEFT,
        )

    # 底部标签
    bindings = slide_data.get("source_bindings", [])
    _add_textbox(
        slide,
        Cm(3.0),
        Cm(11.5),
        Cm(27.0),
        Cm(1.0),
        f"GroundedDeck — {len(bindings)} grounded sources",
        font_size=12,
        color=COLOR_MUTED,
        alignment=PP_ALIGN.LEFT,
    )

    # topic overview（如果有）
    topics = []
    for ve in slide_data.get("visual_elements", []):
        if ve.get("type") == "topic-overview" and "topics" in ve:
            topics = ve["topics"]
            break
    if topics:
        topic_text = "  |  ".join(topics)
        _add_textbox(
            slide,
            Cm(3.0),
            Cm(13.0),
            Cm(27.0),
            Cm(1.5),
            topic_text,
            font_size=13,
            color=COLOR_BODY,
            alignment=PP_ALIGN.LEFT,
        )


def _render_summary(slide, slide_data: dict) -> None:
    """渲染 summary 布局：决策骨架，要点列表 + 来源标注 + claim-source 映射。"""
    _add_title_block(slide, slide_data["title"], slide_data.get("goal", ""))

    key_points = slide_data.get("key_points", [])
    if key_points:
        _add_bullet_list(
            slide,
            MARGIN_LEFT,
            CONTENT_TOP,
            CONTENT_WIDTH,
            Cm(10.0),
            key_points,
        )

    # claim-source-map 可视化（如果有）
    claim_entries = []
    for ve in slide_data.get("visual_elements", []):
        if ve.get("type") == "claim-source-map" and "entries" in ve:
            claim_entries = ve["entries"]
            break
    if claim_entries and not key_points:
        # 如果 key_points 为空但有 claim_entries，作为回退渲染
        items = [f"{e.get('claim', '')} [{e.get('source_binding', '')}]" for e in claim_entries]
        _add_bullet_list(
            slide,
            MARGIN_LEFT,
            CONTENT_TOP,
            CONTENT_WIDTH,
            Cm(10.0),
            items,
        )

    # 底部 source bindings 标签
    bindings = slide_data.get("source_bindings", [])
    if bindings:
        _add_textbox(
            slide,
            MARGIN_LEFT,
            Cm(16.5),
            CONTENT_WIDTH,
            Cm(1.0),
            f"Sources: {', '.join(bindings)}",
            font_size=10,
            color=COLOR_MUTED,
            alignment=PP_ALIGN.LEFT,
        )


def _render_timeline(slide, slide_data: dict) -> None:
    """渲染 timeline 布局：装饰线 + 原生表格（里程碑行 + 事件行）。

    使用 python-pptx 原生 Table 对象展示时间线数据，
    确保用户可以在 PowerPoint 中直接编辑里程碑和事件描述。
    保留水平装饰线作为视觉锚点。
    """
    _add_title_block(slide, slide_data["title"], slide_data.get("goal", ""))

    # 提取 milestones 和 events
    milestones = []
    events = []
    for ve in slide_data.get("visual_elements", []):
        if ve.get("type") == "timeline" and "milestones" in ve:
            milestones = ve["milestones"]
            events = ve.get("events", [])
            break

    if not milestones:
        milestones = ["Start", "Middle", "End"]

    n = len(milestones)

    # 装饰性时间线主轴
    from pptx.enum.shapes import MSO_SHAPE

    line_y = Cm(5.5)
    line_left = Cm(3.0)
    line_width = Cm(27.0)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, line_left, line_y, line_width, Cm(0.15),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_TIMELINE_ACCENT
    line.line.fill.background()

    # 装饰性圆形节点
    spacing = int(line_width) // max(n, 1)
    for i in range(n):
        cx = int(line_left) + spacing * i + spacing // 2
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(cx - int(Cm(0.4))),
            line_y - Cm(0.4),
            Cm(0.8),
            Cm(0.95),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = COLOR_TIMELINE_ACCENT
        dot.line.fill.background()

    # 原生表格：2 行 x n 列（里程碑行 + 事件行）
    table_cols = min(n, 6)
    table_rows = 2
    table_width = CONTENT_WIDTH
    col_width = Emu(int(table_width) // table_cols)
    table_top = Cm(7.0)
    table_height = Cm(7.0)

    table_shape = slide.shapes.add_table(
        table_rows, table_cols,
        MARGIN_LEFT, table_top,
        table_width, table_height,
    )
    table = table_shape.table

    for ci in range(table_cols):
        table.columns[ci].width = col_width

    # 第一行：里程碑标签（绿色强调，居中）
    for ci in range(table_cols):
        cell = table.cell(0, ci)
        cell.text = ""
        _fill_cell(cell, COLOR_TIMELINE_ACCENT)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(milestones[ci]) if ci < n else ""
        _set_font(run, size=18, bold=True, color=COLOR_WHITE, name=FONT_TITLE)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 第二行：事件描述（浅色背景）
    for ci in range(table_cols):
        cell = table.cell(1, ci)
        cell.text = ""
        bg = RGBColor(0xEC, 0xFD, 0xF5) if ci % 2 == 0 else COLOR_WHITE
        _fill_cell(cell, bg)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        if ci < len(events):
            event_text = str(events[ci])
            # 去掉前缀年份（如 "2022: "）以避免重复
            if ":" in event_text:
                event_text = event_text.split(":", 1)[1].strip()
            run.text = event_text
        else:
            run.text = ""
        _set_font(run, size=12, color=COLOR_BODY)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # key_points
    key_points = slide_data.get("key_points", [])
    if key_points:
        kp_top = Emu(int(table_top) + int(table_height) + int(Cm(0.8)))
        _add_textbox(
            slide,
            MARGIN_LEFT,
            kp_top,
            CONTENT_WIDTH,
            Cm(3.0),
            key_points[0],
            font_size=14,
            color=COLOR_BODY,
            alignment=PP_ALIGN.CENTER,
        )


def _render_comparison(slide, slide_data: dict) -> None:
    """渲染 comparison 布局：原生表格对比 + 标题行着色。

    使用 python-pptx 原生 Table 对象，确保输出在 PowerPoint 中
    完全可编辑（可调整列宽、添加行、修改单元格样式）。
    """
    _add_title_block(slide, slide_data["title"], slide_data.get("goal", ""))

    # 提取列名和列要点
    columns = []
    column_points: dict[str, list[str]] = {}
    for ve in slide_data.get("visual_elements", []):
        if ve.get("type") == "comparison-columns" and "columns" in ve:
            columns = ve["columns"]
            column_points = ve.get("column_points", {})
            break

    if len(columns) < 2:
        columns = ["Option A", "Option B"]

    left_points = column_points.get(columns[0], [])[:4]
    right_points = column_points.get(columns[1], [])[:4]
    max_rows = max(len(left_points), len(right_points), 1)

    # 创建原生表格：(max_rows + 1) 行 x 2 列（+1 为标题行）
    table_rows = max_rows + 1
    table_cols = 2
    table_width = CONTENT_WIDTH
    row_height = Cm(1.8)
    table_height = Emu(int(row_height) * table_rows)

    table_shape = slide.shapes.add_table(
        table_rows, table_cols,
        MARGIN_LEFT, CONTENT_TOP,
        table_width, table_height,
    )
    table = table_shape.table

    # 设置列宽（均分）
    col_width = Emu(int(table_width) // 2)
    table.columns[0].width = col_width
    table.columns[1].width = col_width

    # 标题行
    header_colors = [COLOR_COMPARISON_LEFT, COLOR_COMPARISON_RIGHT]
    for ci, col_name in enumerate(columns[:2]):
        cell = table.cell(0, ci)
        cell.text = ""
        _fill_cell(cell, header_colors[ci])
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(col_name)
        _set_font(run, size=18, bold=True, color=COLOR_WHITE, name=FONT_TITLE)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 数据行
    for ri in range(max_rows):
        for ci, points in enumerate([left_points, right_points]):
            cell = table.cell(ri + 1, ci)
            cell.text = ""
            # 交替行背景
            bg = RGBColor(0xEF, 0xF6, 0xFF) if ci == 0 else RGBColor(0xFE, 0xF2, 0xF2)
            if ri % 2 == 1:
                bg = COLOR_WHITE
            _fill_cell(cell, bg)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = f"  {points[ri]}" if ri < len(points) else ""
            _set_font(run, size=14, color=COLOR_BODY)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # key_points 放在表格下方
    key_points = slide_data.get("key_points", [])
    if key_points:
        kp_top = Emu(int(CONTENT_TOP) + int(table_height) + int(Cm(1.0)))
        _add_textbox(
            slide,
            MARGIN_LEFT,
            kp_top,
            CONTENT_WIDTH,
            Cm(2.0),
            key_points[0],
            font_size=14,
            color=COLOR_BODY,
            alignment=PP_ALIGN.CENTER,
        )


def _render_process(slide, slide_data: dict) -> None:
    """渲染 process 布局：原生表格展示步骤 + 装饰性箭头连接。

    使用 python-pptx 原生 Table 对象展示流程步骤，
    确保用户可以在 PowerPoint 中直接编辑步骤内容。
    保留装饰性箭头作为视觉连接。
    """
    _add_title_block(slide, slide_data["title"], slide_data.get("goal", ""))

    # 提取步骤数和步骤标签
    step_count = 1
    ve_step_labels: list[str] = []
    for ve in slide_data.get("visual_elements", []):
        if ve.get("type") == "process-flow":
            step_count = ve.get("steps", 1)
            ve_step_labels = ve.get("step_labels", [])
            break

    # 优先使用 visual_elements 中的 step_labels，回退到 key_points
    if ve_step_labels:
        step_labels = ve_step_labels[:step_count]
    else:
        key_points = slide_data.get("key_points", [])
        step_labels = key_points[:step_count] if key_points else [f"Step {i + 1}" for i in range(max(step_count, 1))]
        # 如果只有一个 key_point 且 steps > 1，拆分成句号分隔的片段
        if len(step_labels) == 1 and step_count > 1:
            parts = step_labels[0].replace("，", "、").split("、")
            step_labels = parts[:step_count]

    n = len(step_labels)
    if n == 0:
        n = 1
        step_labels = [""]

    # 原生表格：2 行 x n 列（序号行 + 内容行）
    table_cols = min(n, 6)
    table_rows = 2
    table_width = CONTENT_WIDTH
    col_width = Emu(int(table_width) // table_cols)
    table_top = Cm(6.0)
    table_height = Cm(7.0)

    table_shape = slide.shapes.add_table(
        table_rows, table_cols,
        MARGIN_LEFT, table_top,
        table_width, table_height,
    )
    table = table_shape.table

    for ci in range(table_cols):
        table.columns[ci].width = col_width

    # 第一行：步骤序号（紫色强调）
    for ci in range(table_cols):
        cell = table.cell(0, ci)
        cell.text = ""
        # 最后一步用紫色实底，其他用浅色
        is_last = (ci == table_cols - 1)
        bg = COLOR_PROCESS_ACCENT if is_last else COLOR_LIGHT_BG
        _fill_cell(cell, bg)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"Step {ci + 1}"
        text_color = COLOR_WHITE if is_last else COLOR_PROCESS_ACCENT
        _set_font(run, size=14, bold=True, color=text_color, name=FONT_TITLE)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 第二行：步骤内容
    for ci in range(table_cols):
        cell = table.cell(1, ci)
        cell.text = ""
        is_last = (ci == table_cols - 1)
        bg = COLOR_PROCESS_ACCENT if is_last else COLOR_WHITE
        _fill_cell(cell, bg)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = step_labels[ci] if ci < len(step_labels) else ""
        text_color = COLOR_WHITE if is_last else COLOR_BODY
        _set_font(run, size=12, color=text_color)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 装饰性箭头（在表格下方，连接各列）
    from pptx.enum.shapes import MSO_SHAPE

    arrow_y = Emu(int(table_top) + int(table_height) + int(Cm(0.3)))
    arrow_height = Cm(0.6)
    for ci in range(table_cols - 1):
        arrow_x = Emu(int(MARGIN_LEFT) + int(col_width) * ci + int(col_width) - int(Cm(0.5)))
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            arrow_x,
            arrow_y,
            Cm(1.0),
            arrow_height,
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLOR_PROCESS_ACCENT
        arrow.line.fill.background()

    # key_points（放在箭头下方）
    key_points = slide_data.get("key_points", [])
    if key_points:
        kp_top = Emu(int(arrow_y) + int(arrow_height) + int(Cm(0.8)))
        _add_textbox(
            slide,
            MARGIN_LEFT,
            kp_top,
            CONTENT_WIDTH,
            Cm(3.0),
            key_points[0],
            font_size=14,
            color=COLOR_BODY,
            alignment=PP_ALIGN.CENTER,
        )


def _render_chart(slide, slide_data: dict) -> None:
    """渲染 chart 布局：原生表格展示指标 + 关键数字 + 上下文标签。

    使用 python-pptx 原生 Table 对象展示指标数据，
    确保用户可以在 PowerPoint 中直接编辑数值和标签。
    当指标数量 <= 4 时使用横向表格，否则回退到纵向表格。
    """
    _add_title_block(slide, slide_data["title"], slide_data.get("goal", ""))

    # 提取指标和标签
    metrics = []
    metric_labels = []
    for ve in slide_data.get("visual_elements", []):
        if ve.get("type") == "metric-cards" and "metrics" in ve:
            metrics = ve["metrics"]
            metric_labels = ve.get("labels", [])
            break

    if not metrics:
        metrics = ["—"]

    n = len(metrics)

    # 使用原生表格展示指标：2 行 x n 列（标签行 + 数值行）
    table_rows = 2
    table_cols = min(n, 6)  # 最多 6 列
    table_width = CONTENT_WIDTH
    col_width = Emu(int(table_width) // table_cols)
    table_height = Cm(6.0)

    table_shape = slide.shapes.add_table(
        table_rows, table_cols,
        MARGIN_LEFT, Cm(7.0),
        table_width, table_height,
    )
    table = table_shape.table

    # 设置列宽
    for ci in range(table_cols):
        table.columns[ci].width = col_width

    # 第一行：指标数值（大字体，橙色强调）
    for ci in range(table_cols):
        cell = table.cell(0, ci)
        cell.text = ""
        _fill_cell(cell, COLOR_LIGHT_BG)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(metrics[ci]) if ci < n else "—"
        _set_font(run, size=32, bold=True, color=COLOR_CHART_ACCENT, name=FONT_TITLE)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 第二行：指标标签（小字体，灰色）
    for ci in range(table_cols):
        cell = table.cell(1, ci)
        cell.text = ""
        _fill_cell(cell, COLOR_WHITE)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        label_text = metric_labels[ci] if ci < len(metric_labels) else f"Metric {ci + 1}"
        run.text = label_text
        _set_font(run, size=12, color=COLOR_MUTED)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # key_points
    key_points = slide_data.get("key_points", [])
    if key_points:
        _add_textbox(
            slide,
            MARGIN_LEFT,
            Cm(14.5),
            CONTENT_WIDTH,
            Cm(3.0),
            key_points[0],
            font_size=14,
            color=COLOR_BODY,
            alignment=PP_ALIGN.CENTER,
        )


def _render_section(slide, slide_data: dict) -> None:
    """渲染 section 布局：章节分隔页，带装饰性分隔线。"""
    from pptx.enum.shapes import MSO_SHAPE

    # 背景色块
    _add_rounded_rect(slide, Cm(0), Cm(6.0), SLIDE_WIDTH, Cm(7.0), fill_color=COLOR_PRIMARY)

    # 上方装饰线
    top_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(6.0), Cm(6.8), Cm(21.0), Cm(0.08),
    )
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = COLOR_ACCENT
    top_line.line.fill.background()

    # 标题
    _add_textbox(
        slide,
        Cm(3.0),
        Cm(7.5),
        Cm(27.0),
        Cm(3.0),
        slide_data["title"],
        font_size=32,
        bold=True,
        color=COLOR_WHITE,
        alignment=PP_ALIGN.CENTER,
        font_name=FONT_TITLE,
    )

    # 目标描述
    goal = slide_data.get("goal", "")
    if goal:
        _add_textbox(
            slide,
            Cm(3.0),
            Cm(10.5),
            Cm(27.0),
            Cm(1.5),
            goal,
            font_size=16,
            color=RGBColor(0xBF, 0xDB, 0xFE),
            alignment=PP_ALIGN.CENTER,
        )

    # 下方装饰线
    bottom_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(6.0), Cm(12.2), Cm(21.0), Cm(0.08),
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = COLOR_ACCENT
    bottom_line.line.fill.background()


def _fill_cell(cell, color: RGBColor) -> None:
    """设置表格单元格的背景填充色。"""
    from lxml import etree
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    # 移除已有的 solidFill
    for sf in tcPr.findall("a:solidFill", nsmap):
        tcPr.remove(sf)
    solid_fill = etree.SubElement(
        tcPr,
        "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill",
    )
    srgb = etree.SubElement(
        solid_fill,
        "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr",
    )
    srgb.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")


# 布局路由表
_LAYOUT_RENDERERS: dict[str, Any] = {
    "cover": _render_cover,
    "summary": _render_summary,
    "timeline": _render_timeline,
    "comparison": _render_comparison,
    "process": _render_process,
    "chart": _render_chart,
    "section": _render_section,
}


# ================================================================
# 公开 API
# ================================================================


def render_slide_spec_to_pptx(
    slide_spec: dict,
    output_path: str | Path,
) -> Path:
    """将 slide spec 渲染为可编辑的 .pptx 文件。

    参数：
        slide_spec: 符合 slide-spec schema 的 dict
        output_path: 输出 .pptx 文件路径

    返回：
        输出文件的 Path 对象

    异常：
        ValueError: slide_spec 缺少必要字段
        KeyError: slide 使用了不支持的 layout_type
    """
    output_path = Path(output_path)

    if "slides" not in slide_spec:
        raise ValueError("slide_spec must contain a 'slides' list")
    if not slide_spec["slides"]:
        raise ValueError("slide_spec must contain at least one slide")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 使用空白布局
    blank_layout = prs.slide_layouts[6]  # 空白 layout

    for slide_data in slide_spec["slides"]:
        layout_type = slide_data.get("layout_type", "summary")
        renderer = _LAYOUT_RENDERERS.get(layout_type)

        if renderer is None:
            logger.warning(
                "unsupported layout_type '%s' for slide '%s', falling back to summary",
                layout_type,
                slide_data.get("slide_id", "unknown"),
            )
            renderer = _render_summary

        slide = prs.slides.add_slide(blank_layout)
        renderer(slide, slide_data)
        _add_speaker_notes(slide, slide_data)

    # 确保输出目录存在
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    logger.info("rendered %d slides to %s", len(slide_spec["slides"]), output_path)
    return output_path


def get_supported_layouts() -> tuple[str, ...]:
    """返回渲染器支持的所有 layout_type。"""
    return tuple(sorted(_LAYOUT_RENDERERS.keys()))
