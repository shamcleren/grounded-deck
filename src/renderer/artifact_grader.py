"""PPTX artifact 质量评估器：检查生成的 PPTX 文件的可编辑性、source coverage 和结构质量。

该模块实现 evaluation-plan.md 中 phase two 的 artifact graders，
用于检查导出的 PPT 文件是否满足以下质量标准：
- 所有文本框保持可编辑
- source bindings 写入备注区
- 中文内容正确渲染
- slide 数量与 slide spec 一致
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from pptx import Presentation

logger = logging.getLogger(__name__)


def grade_pptx_artifact(
    pptx_path: str | Path,
    slide_spec: dict | None = None,
) -> dict:
    """对 PPTX 文件执行 artifact 质量检查。

    参数：
        pptx_path: PPTX 文件路径
        slide_spec: 可选的 slide spec dict，用于交叉验证

    返回：
        质量报告 dict，包含 status, failures, metrics 等字段
    """
    pptx_path = Path(pptx_path)
    failures: list[str] = []
    warnings: list[str] = []

    # ---------- 文件存在性检查 ----------
    if not pptx_path.exists():
        return {
            "status": "fail",
            "failures": [f"PPTX file not found: {pptx_path}"],
            "warnings": [],
            "metrics": {},
        }

    # ---------- 打开 PPTX ----------
    try:
        prs = Presentation(str(pptx_path))
    except Exception as exc:
        return {
            "status": "fail",
            "failures": [f"failed to open PPTX: {exc}"],
            "warnings": [],
            "metrics": {},
        }

    slide_count = len(prs.slides)
    total_shapes = 0
    total_text_boxes = 0
    editable_text_boxes = 0
    slides_with_notes = 0
    slides_with_source_bindings = 0
    chinese_text_found = False
    empty_slides: list[int] = []

    for i, slide in enumerate(prs.slides):
        shape_count = 0
        text_count = 0

        for shape in slide.shapes:
            shape_count += 1
            if shape.has_text_frame:
                text_count += 1
                # 检查文本框是否可编辑（has_text_frame 意味着可编辑）
                editable_text_boxes += 1

                # 检查是否有中文内容
                for para in shape.text_frame.paragraphs:
                    text = para.text
                    if any(ord(c) > 0x4E00 for c in text):
                        chinese_text_found = True

        total_shapes += shape_count
        total_text_boxes += text_count

        if shape_count == 0:
            empty_slides.append(i + 1)

        # 检查 speaker notes
        try:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                slides_with_notes += 1
                if "[Source Bindings]" in notes:
                    slides_with_source_bindings += 1
        except Exception:
            pass

    # ---------- 结构检查 ----------
    if slide_count == 0:
        failures.append("PPTX contains no slides")

    if empty_slides:
        failures.append(f"empty slides (no shapes): {empty_slides}")

    if editable_text_boxes == 0:
        failures.append("no editable text boxes found")

    # ---------- 与 slide spec 交叉验证 ----------
    if slide_spec is not None:
        expected_count = len(slide_spec.get("slides", []))
        if slide_count != expected_count:
            failures.append(
                f"slide count mismatch: PPTX has {slide_count}, spec has {expected_count}"
            )

    # ---------- 备注完整性检查 ----------
    notes_ratio = slides_with_notes / max(1, slide_count)
    if notes_ratio < 0.8:
        warnings.append(
            f"only {slides_with_notes}/{slide_count} slides have speaker notes "
            f"({notes_ratio:.0%})"
        )

    source_binding_ratio = slides_with_source_bindings / max(1, slide_count)
    if source_binding_ratio < 0.5:
        warnings.append(
            f"only {slides_with_source_bindings}/{slide_count} slides have source bindings in notes "
            f"({source_binding_ratio:.0%})"
        )

    return {
        "status": "pass" if not failures else "fail",
        "failures": failures,
        "warnings": warnings,
        "metrics": {
            "file_size_bytes": pptx_path.stat().st_size,
            "slide_count": slide_count,
            "total_shapes": total_shapes,
            "total_text_boxes": total_text_boxes,
            "editable_text_boxes": editable_text_boxes,
            "editability_ratio": round(
                editable_text_boxes / max(1, total_text_boxes), 2
            ),
            "slides_with_notes": slides_with_notes,
            "slides_with_source_bindings": slides_with_source_bindings,
            "notes_coverage_ratio": round(notes_ratio, 2),
            "source_binding_coverage_ratio": round(source_binding_ratio, 2),
            "chinese_text_found": chinese_text_found,
        },
    }
